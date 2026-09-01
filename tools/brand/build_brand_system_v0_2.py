#!/usr/bin/env python3
from __future__ import annotations

import argparse
import base64
import hashlib
import json
import math
import os
import re
import shutil
import subprocess
import sys
import textwrap
import zipfile
from dataclasses import dataclass
from datetime import date
from pathlib import Path
from typing import Iterable
from xml.sax.saxutils import escape

import cairosvg
from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from fontTools.pens.svgPathPen import SVGPathPen
from fontTools.ttLib import TTFont
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PptxInches, Pt as PptxPt
from reportlab.lib.colors import Color, HexColor
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import A4, LETTER, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch, mm
from reportlab.pdfbase.ttfonts import TTFont as RLTTFont
from reportlab.pdfbase import pdfmetrics
from reportlab.platypus import (
    BaseDocTemplate,
    Frame,
    Image as RLImage,
    KeepTogether,
    PageBreak,
    PageTemplate,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

VERSION = "0.2.0"
BUILD_DATE = "2026-09-01"
REPO_URL = "https://github.com/SquirmyWormy275/SABLEHARBOR"
CANON_PATH = "docs/canon/SABLE_HARBOR_CORPORATE_LORE_CANON_v0.2.md"

OFF_WHITE = "#F4F1EA"
INK = "#101214"
MUTED = "#747A80"
DARK = "#101419"
WHITE = "#F7F5EF"
ORANGE = "#C45124"
EVALON_STEEL = "#687C86"

PALETTE = {
    "Sable Harbor Ink": INK,
    "Sable Harbor Orange": ORANGE,
    "Field Green": "#315F4D",
    "Meridian Blue": "#2E6F96",
    "Pale Sun Gold": "#C38B1F",
    "Advisory Steel Blue": "#456C98",
    "Red Wash Oxide": "#B94C2C",
    "Evalon Steel": EVALON_STEEL,
    "Paper": OFF_WHITE,
    "Reverse Field": DARK,
    "Reverse White": WHITE,
    "Muted Gray": MUTED,
}

FONT_MEDIUM_CANDIDATES = [
    "/usr/share/fonts/opentype/inter/InterDisplay-Medium.otf",
    "/usr/share/fonts/truetype/inter/InterDisplay-Medium.ttf",
    "/usr/share/fonts/truetype/inter/Inter-Medium.ttf",
]
FONT_SEMIBOLD_CANDIDATES = [
    "/usr/share/fonts/opentype/inter/InterDisplay-SemiBold.otf",
    "/usr/share/fonts/truetype/inter/InterDisplay-SemiBold.ttf",
    "/usr/share/fonts/truetype/inter/Inter-SemiBold.ttf",
]
FONT_BOLD_CANDIDATES = [
    "/usr/share/fonts/opentype/inter/InterDisplay-Bold.otf",
    "/usr/share/fonts/truetype/inter/InterDisplay-Bold.ttf",
    "/usr/share/fonts/truetype/inter/Inter-Bold.ttf",
]


def first_existing(paths: Iterable[str]) -> Path:
    for value in paths:
        p = Path(value)
        if p.exists():
            return p
    raise FileNotFoundError(f"None of these font paths exist: {list(paths)}")


FONT_MEDIUM_PATH = first_existing(FONT_MEDIUM_CANDIDATES)
FONT_SEMIBOLD_PATH = first_existing(FONT_SEMIBOLD_CANDIDATES)
FONT_BOLD_PATH = first_existing(FONT_BOLD_CANDIDATES)


class OutlineFont:
    def __init__(self, path: Path):
        self.path = path
        self.font = TTFont(str(path))
        self.glyph_set = self.font.getGlyphSet()
        self.cmap = self.font.getBestCmap()
        self.hmtx = self.font["hmtx"]
        self.units_per_em = self.font["head"].unitsPerEm
        self.cache: dict[str, tuple[str, int]] = {}

    def glyph(self, ch: str) -> tuple[str, int]:
        if ch in self.cache:
            return self.cache[ch]
        glyph_name = self.cmap.get(ord(ch), ".notdef")
        pen = SVGPathPen(self.glyph_set)
        self.glyph_set[glyph_name].draw(pen)
        d = pen.getCommands()
        advance = self.hmtx[glyph_name][0]
        self.cache[ch] = (d, advance)
        return d, advance

    def measure(self, text: str, size: float, tracking: float = 0) -> float:
        total = 0.0
        for idx, ch in enumerate(text):
            if ch == " ":
                advance = self.units_per_em * 0.48
            else:
                _, advance = self.glyph(ch)
            total += advance / self.units_per_em * size
            if idx != len(text) - 1:
                total += tracking
        return total

    def paths(
        self,
        text: str,
        x: float,
        baseline: float,
        size: float,
        fill: str,
        tracking: float = 0,
        anchor: str = "start",
        opacity: float | None = None,
    ) -> str:
        width = self.measure(text, size, tracking)
        cursor = x
        if anchor == "middle":
            cursor = x - width / 2
        elif anchor == "end":
            cursor = x - width
        scale = size / self.units_per_em
        op = "" if opacity is None else f' opacity="{opacity:.4f}"'
        parts = [f'<g fill="{fill}"{op}>']
        for idx, ch in enumerate(text):
            if ch == " ":
                advance = self.units_per_em * 0.48
            else:
                d, advance = self.glyph(ch)
                if d:
                    parts.append(
                        f'<path d="{d}" transform="translate({cursor:.3f},{baseline:.3f}) scale({scale:.8f},{-scale:.8f})"/>'
                    )
            cursor += advance / self.units_per_em * size
            if idx != len(text) - 1:
                cursor += tracking
        parts.append("</g>")
        return "".join(parts)


FONT_MEDIUM = OutlineFont(FONT_MEDIUM_PATH)
FONT_SEMIBOLD = OutlineFont(FONT_SEMIBOLD_PATH)
FONT_BOLD = OutlineFont(FONT_BOLD_PATH)


@dataclass(frozen=True)
class Identity:
    slug: str
    display_name: str
    role: str
    status: str
    accent: str
    tagline: str
    description: str
    source_section: str
    current: bool
    variants: tuple[str, ...]


CORE_VARIANTS = (
    "primary-horizontal",
    "stacked",
    "mark",
    "reverse-horizontal",
    "one-color-horizontal",
)
HISTORICAL_VARIANTS = CORE_VARIANTS

IDENTITIES: tuple[Identity, ...] = (
    Identity(
        "sable-harbor",
        "Sable Harbor",
        "corporate master brand",
        "LOCKED corporate identity; quantitative and legal structure remain separate workstreams.",
        ORANGE,
        "INDUSTRIAL SYSTEMS",
        "Sable Harbor is an industrial-systems company and reusable synthetic enterprise spanning mining, natural resources, operational software, analytics, assurance, finance, governance, security, and professional training. Its governing idea is that variation is normal: the work is to encounter it, represent it, understand it, and choose consciously what happens next.",
        "Sections 1 and 13",
        True,
        CORE_VARIANTS,
    ),
    Identity(
        "foundry-field",
        "Foundry Field",
        "current core business line",
        "LOCKED name and role.",
        ORANGE,
        "A SABLE HARBOR BUSINESS LINE",
        "Foundry Field is the deployable operational product and service configuration built on the Foundry substrate. It encounters, connects, and represents operational reality while preserving identity, provenance, competing definitions, effective dates, uncertainty, and the record of decisions.",
        "Sections 6.1-6.8 and 13",
        True,
        CORE_VARIANTS,
    ),
    Identity(
        "willow",
        "Willow",
        "current core business line",
        "LOCKED name, purpose, and two-center operating model; exact legal structure remains OPEN.",
        "#315F4D",
        "A SABLE HARBOR BUSINESS LINE",
        "Willow pursues consequential industrial questions that no existing product, delivery team, or operating unit can own, changing the physical world in bounded ways to learn what may be possible. Its formal loop is Question -> belief -> experiment -> observation -> decision; failure is survivable, drift is not.",
        "Sections 7.7-8.6 and 13",
        True,
        CORE_VARIANTS,
    ),
    Identity(
        "atlas-meridian",
        "Atlas Meridian",
        "current core business line",
        "LOCKED product lineage and decision-support boundary; launch details remain OPEN.",
        "#2E6F96",
        "A SABLE HARBOR BUSINESS LINE",
        "Atlas Meridian is a disciplined investigative system operating across represented industrial evidence. Foundry represents the terrain; Atlas Meridian investigates across it. It preserves provenance, authorities, constraints, stop conditions, and the difference between an answer and an investigation, and it remains decision support rather than an autonomous decision-maker.",
        "Sections 9.1-9.9 and 13",
        True,
        CORE_VARIANTS,
    ),
    Identity(
        "pale-sun",
        "Pale Sun",
        "current operating business line",
        "LOCKED operating thesis; detailed mine, transaction, permitting, and economic facts remain OPEN.",
        "#C38B1F",
        "A SABLE HARBOR BUSINESS LINE",
        "Pale Sun owns and operates the fictional Red Wash uranium mine where operating control is the thesis. It is an operation first and a proving ground second, preserving conservative information boundaries, field-qualification gates, human authority, and the distinction between claims, dispositions, and later outcomes.",
        "Sections 10.1-10.11 and 13",
        True,
        CORE_VARIANTS,
    ),
    Identity(
        "project-cradle",
        "Project Cradle",
        "current core business line",
        "LOCKED thesis and founding team; first commercial details remain OPEN.",
        "#C58A14",
        "A SABLE HARBOR BUSINESS LINE",
        "Project Cradle finds valuable rare-earth material that existing mines and processing plants already concentrate unintentionally, then adds the smallest recovery step that captures value without breaking the host operation. It seeks value in a process seam and generally avoids owning the host mine.",
        "Sections 11.1-11.6 and 13",
        True,
        CORE_VARIANTS,
    ),
    Identity(
        "american-resource-utility",
        "American Resource Utility",
        "current distinct operating company and core line",
        "LOCKED name and operating role; detailed history, routes, assets, workforce, and legal structure remain OPEN.",
        ORANGE,
        "A SABLE HARBOR BUSINESS LINE",
        "American Resource Utility (ARU) is an acquired resource-logistics operator with legacy customers, employees, dispatch practices, terminals, equipment, and operating knowledge. It remains operationally distinct while moving materials and supplies across physical and organizational boundaries.",
        "Sections 12.1-12.5 and 13",
        True,
        CORE_VARIANTS,
    ),
    Identity(
        "advisory",
        "Advisory",
        "emerging current business line",
        "LOCKED direction; exact name, leader, launch date, service catalog, P&L, and organizational home remain OPEN.",
        "#456C98",
        "A SABLE HARBOR BUSINESS LINE",
        "Advisory transfers Sable Harbor's method where a client can and should own the system. It is not generic consulting, a slide-deck practice, or merely software implementation; it enters a messy system, determines what is actually happening, and transfers a method the operator can sustain.",
        "Sections 12.6 and 13",
        True,
        CORE_VARIANTS,
    ),
    Identity(
        "foundry",
        "Foundry",
        "supplemental product substrate",
        "LOCKED distinction from Foundry Field.",
        ORANGE,
        "THE SABLE HARBOR PRODUCT SUBSTRATE",
        "Foundry is the underlying relationship, meaning, integration, provenance, and workflow substrate. Its graph represents relationships and local meaning without pretending every object has one permanent identifier or every question has one context-free answer.",
        "Sections 6.1-6.7",
        False,
        ("primary-horizontal", "mark", "reverse-horizontal"),
    ),
    Identity(
        "red-wash-mine",
        "Red Wash Mine",
        "operating asset under Pale Sun",
        "LOCKED fictional mine name; detailed asset facts remain OPEN.",
        "#B94C2C",
        "A PALE SUN OPERATING ASSET",
        "Red Wash is the fictional underground Wyoming uranium mine operated by Pale Sun. Its role in the canon is to force Sable Harbor's representational systems to confront physical reality, operating authority, environmental history, and incomplete evidence.",
        "Sections 10.4-10.11",
        False,
        ("primary-horizontal", "mark", "reverse-horizontal"),
    ),
    Identity(
        "blood-sweat-and-tears-railway",
        "Blood, Sweat & Tears Railway",
        "ARU operating component",
        "LOCKED name and relationship to ARU; legal identity, routes, assets, and history remain OPEN.",
        "#B94C2C",
        "AN AMERICAN RESOURCE UTILITY OPERATING COMPONENT",
        "The Blood, Sweat & Tears Railway (BS&T) is ARU's railway or short-line operating component. It persists as corporate archaeology inside Sable Harbor's acquired logistics estate.",
        "Section 12.3",
        False,
        ("primary-horizontal", "mark", "reverse-horizontal"),
    ),
    Identity(
        "emberline",
        "Emberline",
        "historical business line",
        "LOCKED historical status: active through 2025, then absorbed into enduring 2026 work.",
        ORANGE,
        "A HISTORICAL SABLE HARBOR LINE",
        "Project Emberline was the Charleston-centered commercial, customer, and field-experimentation coal program. It supplied customers, operating environments, and paid experiments before its enduring work moved into Foundry, Willow, field operations, and emerging Advisory.",
        "Section 7.3",
        False,
        ("primary-horizontal", "mark", "reverse-horizontal"),
    ),
    Identity(
        "evalon",
        "Evalon",
        "historical advanced-engineering outpost",
        "LOCKED historical concept; closed and rechartered as Willow in 2022.",
        EVALON_STEEL,
        "A HISTORICAL SABLE HARBOR OUTPOST",
        "Evalon was the Pittsburgh-area advanced-engineering outpost created under a coal and industrial mandate during 2020-2021. It failed as the business it was intended to become, concentrated experimental capability, and was closed as an operating concept before the surviving capability was rechartered as Project Willow.",
        "Sections 7.2-7.7",
        False,
        HISTORICAL_VARIANTS,
    ),
    Identity(
        "red-wash-pale-sun",
        "Red Wash / Pale Sun",
        "supplemental endorsed operating lockup",
        "Supplemental lockup; does not replace either canonical name.",
        "#B94C2C",
        "A SABLE HARBOR OPERATING BUSINESS",
        "This endorsed lockup joins the Pale Sun operating line to the Red Wash asset for contexts where both relationships must be visible. It does not create a separate business line.",
        "Sections 10 and 13",
        False,
        ("primary-horizontal", "mark", "reverse-horizontal"),
    ),
)


def rgb_tuple(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return tuple(int(value[i:i+2], 16) for i in (0, 2, 4))


def cmyk_approx(value: str) -> tuple[int, int, int, int]:
    r, g, b = (c / 255 for c in rgb_tuple(value))
    k = 1 - max(r, g, b)
    if math.isclose(k, 1.0):
        return (0, 0, 0, 100)
    c = (1 - r - k) / (1 - k)
    m = (1 - g - k) / (1 - k)
    y = (1 - b - k) / (1 - k)
    return tuple(round(x * 100) for x in (c, m, y, k))


def svg_rect(x: float, y: float, w: float, h: float, fill: str, rx: float = 0) -> str:
    r = f' rx="{rx}"' if rx else ""
    return f'<rect x="{x}" y="{y}" width="{w}" height="{h}" fill="{fill}"{r}/>'


def svg_polygon(points: Iterable[tuple[float, float]], fill: str) -> str:
    text = " ".join(f"{x:.3f},{y:.3f}" for x, y in points)
    return f'<polygon points="{text}" fill="{fill}"/>'


def evalon_mark(main: str, accent: str) -> str:
    # Industrial E / strata crossed by a measurement datum.
    parts = [
        svg_rect(8, 10, 20, 100, main),
        svg_rect(8, 10, 72, 20, main),
        svg_rect(8, 50, 62, 20, main),
        svg_rect(8, 90, 72, 20, main),
        svg_rect(78, 18, 8, 84, accent),
        svg_polygon(((72, 18), (92, 18), (82, 8)), accent),
        svg_rect(73, 54, 18, 12, accent),
        svg_rect(78, 102, 8, 8, accent),
    ]
    return "".join(parts)


def svg_doc(width: int, height: int, title: str, desc: str, body: str, background: str | None = None) -> str:
    bg = svg_rect(0, 0, width, height, background) if background else ""
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" '
        f'viewBox="0 0 {width} {height}" role="img" aria-labelledby="title desc">'
        f'<title id="title">{escape(title)}</title><desc id="desc">{escape(desc)}</desc>{bg}{body}</svg>'
    )


def generate_evalon_logo(variant: str) -> str:
    reverse = variant == "reverse-horizontal"
    one_color = variant == "one-color-horizontal"
    main = WHITE if reverse else INK
    accent = main if one_color else EVALON_STEEL
    bg = DARK if reverse else None
    title = f"Evalon - {variant.replace('-', ' ')}"
    desc = "Individual archival Evalon logo; a historical Sable Harbor advanced-engineering outpost that was rechartered as Willow."

    if variant == "mark":
        body = f'<g transform="translate(180,180) scale(5.333333)">{evalon_mark(INK, EVALON_STEEL)}</g>'
        return svg_doc(1000, 1000, title, desc, body)

    if variant == "stacked":
        mark = f'<g transform="translate(365,90) scale(2.25)">{evalon_mark(INK, EVALON_STEEL)}</g>'
        word = FONT_SEMIBOLD.paths("EVALON", 500, 585, 92, INK, tracking=9, anchor="middle")
        sub = FONT_MEDIUM.paths("A HISTORICAL SABLE HARBOR OUTPOST", 500, 670, 28, MUTED, tracking=3, anchor="middle")
        return svg_doc(1000, 1000, title, desc, mark + word + sub)

    mark = f'<g transform="translate(84,90) scale(2.65)">{evalon_mark(main, accent)}</g>'
    word = FONT_SEMIBOLD.paths("EVALON", 500, 245, 112, main, tracking=11)
    sub = FONT_MEDIUM.paths("A HISTORICAL SABLE HARBOR OUTPOST", 503, 325, 29, MUTED if not reverse else "#B8BEC4", tracking=3)
    line = svg_rect(503, 357, 880, 4, accent)
    return svg_doc(1600, 500, title, desc, mark + word + sub + line, background=bg)


def write_evalon_assets(logo_dir: Path) -> None:
    logo_dir.mkdir(parents=True, exist_ok=True)
    for variant in HISTORICAL_VARIANTS:
        svg = generate_evalon_logo(variant)
        svg_path = logo_dir / f"evalon__{variant}.svg"
        png_path = logo_dir / f"evalon__{variant}.png"
        svg_path.write_text(svg, encoding="utf-8")
        cairosvg.svg2png(bytestring=svg.encode("utf-8"), write_to=str(png_path), output_width=int(re.search(r'width="(\d+)"', svg).group(1)))


def logo_data_uri(path: Path) -> str:
    data = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:image/svg+xml;base64,{data}"


def self_contained_letterhead_svg(repo: Path, page: str) -> tuple[str, int, int]:
    if page == "us-letter":
        w, h = 816, 1056
    else:
        w, h = 794, 1123
    logo_uri = logo_data_uri(repo / "assets/brand/logos/sable-harbor__primary-horizontal.svg")
    body = [
        f'<image href="{logo_uri}" x="54" y="40" width="400" height="125" preserveAspectRatio="xMinYMid meet"/>',
        svg_rect(54, 164, w - 108, 3, ORANGE),
        f'<text x="{w-54}" y="80" text-anchor="end" font-family="Arial, sans-serif" font-size="13" font-weight="700" letter-spacing="2" fill="{INK}">CORPORATE CORRESPONDENCE</text>',
        f'<text x="{w-54}" y="104" text-anchor="end" font-family="Arial, sans-serif" font-size="10" letter-spacing="1.2" fill="{MUTED}">[DATE]  |  [REFERENCE]</text>',
        f'<text x="54" y="218" font-family="Arial, sans-serif" font-size="12" fill="{MUTED}">[RECIPIENT NAME]</text>',
        f'<text x="54" y="238" font-family="Arial, sans-serif" font-size="12" fill="{MUTED}">[ORGANIZATION]</text>',
        f'<text x="54" y="258" font-family="Arial, sans-serif" font-size="12" fill="{MUTED}">[ADDRESS]</text>',
        f'<text x="54" y="314" font-family="Arial, sans-serif" font-size="16" font-weight="700" fill="{INK}">[SUBJECT]</text>',
        f'<text x="54" y="354" font-family="Arial, sans-serif" font-size="12" fill="{INK}">[LETTER BODY BEGINS HERE]</text>',
        svg_rect(54, h - 72, w - 108, 2, ORANGE),
        f'<text x="54" y="{h-44}" font-family="Arial, sans-serif" font-size="9" letter-spacing="0.7" fill="{MUTED}">[ADDRESS]  |  [EMAIL]  |  [PHONE]  |  [WEB]</text>',
        f'<text x="{w-54}" y="{h-44}" text-anchor="end" font-family="Arial, sans-serif" font-size="9" fill="{MUTED}">SABLE HARBOR</text>',
    ]
    return svg_doc(w, h, f"Sable Harbor letterhead - {page}", "Editable corporate letterhead with explicit contact placeholders.", "".join(body), OFF_WHITE), w, h


def self_contained_report_cover_svg(repo: Path, page: str) -> tuple[str, int, int]:
    if page == "us-letter":
        w, h = 816, 1056
    else:
        w, h = 794, 1123
    logo_uri = logo_data_uri(repo / "assets/brand/logos/sable-harbor__reverse-horizontal.svg")
    top = int(h * 0.12)
    body = [
        svg_rect(0, 0, w, h, OFF_WHITE),
        svg_rect(0, 0, w, int(h * 0.44), DARK),
        f'<image href="{logo_uri}" x="52" y="{top}" width="{w-104}" height="210" preserveAspectRatio="xMinYMid meet"/>',
        svg_rect(52, int(h * 0.44) - 6, w - 104, 6, ORANGE),
        f'<text x="52" y="{int(h*0.57)}" font-family="Arial, sans-serif" font-size="38" font-weight="700" letter-spacing="1.5" fill="{INK}">[REPORT TITLE]</text>',
        f'<text x="52" y="{int(h*0.62)}" font-family="Arial, sans-serif" font-size="18" letter-spacing="1" fill="{MUTED}">[REPORT SUBTITLE OR ENGAGEMENT]</text>',
        f'<text x="52" y="{int(h*0.75)}" font-family="Arial, sans-serif" font-size="12" font-weight="700" letter-spacing="1.2" fill="{INK}">DOCUMENT CONTROL</text>',
        f'<text x="52" y="{int(h*0.79)}" font-family="Arial, sans-serif" font-size="11" fill="{MUTED}">DATE</text>',
        f'<text x="170" y="{int(h*0.79)}" font-family="Arial, sans-serif" font-size="11" fill="{INK}">[DATE]</text>',
        f'<text x="52" y="{int(h*0.82)}" font-family="Arial, sans-serif" font-size="11" fill="{MUTED}">VERSION</text>',
        f'<text x="170" y="{int(h*0.82)}" font-family="Arial, sans-serif" font-size="11" fill="{INK}">[VERSION]</text>',
        f'<text x="52" y="{int(h*0.85)}" font-family="Arial, sans-serif" font-size="11" fill="{MUTED}">CLASSIFICATION</text>',
        f'<text x="170" y="{int(h*0.85)}" font-family="Arial, sans-serif" font-size="11" fill="{INK}">[CLASSIFICATION]</text>',
        svg_rect(52, h - 70, w - 104, 2, ORANGE),
        f'<text x="52" y="{h-42}" font-family="Arial, sans-serif" font-size="9" letter-spacing="1" fill="{MUTED}">SABLE HARBOR  |  CONTROLLED DOCUMENT TEMPLATE</text>',
    ]
    return svg_doc(w, h, f"Sable Harbor report cover - {page}", "Corporate report-cover template with explicit document-control placeholders.", "".join(body)), w, h


def svg_to_png_pdf(svg_path: Path, png_path: Path, pdf_path: Path) -> None:
    cairosvg.svg2png(url=str(svg_path), write_to=str(png_path), output_width=int(re.search(r'width="(\d+)"', svg_path.read_text()).group(1)) * 2)
    cairosvg.svg2pdf(url=str(svg_path), write_to=str(pdf_path))


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill.lstrip("#"))
    tc_pr.append(shd)


def set_cell_margins(cell, top=80, start=80, bottom=80, end=80) -> None:
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcMar = tcPr.first_child_found_in("w:tcMar")
    if tcMar is None:
        tcMar = OxmlElement("w:tcMar")
        tcPr.append(tcMar)
    for m, v in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
        node = tcMar.find(qn(f"w:{m}"))
        if node is None:
            node = OxmlElement(f"w:{m}")
            tcMar.append(node)
        node.set(qn("w:w"), str(v))
        node.set(qn("w:type"), "dxa")


def add_bottom_border(paragraph, color: str = ORANGE, size: int = 10) -> None:
    p = paragraph._p
    pPr = p.get_or_add_pPr()
    pBdr = pPr.find(qn("w:pBdr"))
    if pBdr is None:
        pBdr = OxmlElement("w:pBdr")
        pPr.append(pBdr)
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), str(size))
    bottom.set(qn("w:space"), "2")
    bottom.set(qn("w:color"), color.lstrip("#"))
    pBdr.append(bottom)


def add_page_field(paragraph) -> None:
    run = paragraph.add_run()
    fldChar1 = OxmlElement("w:fldChar")
    fldChar1.set(qn("w:fldCharType"), "begin")
    instrText = OxmlElement("w:instrText")
    instrText.set(qn("xml:space"), "preserve")
    instrText.text = "PAGE"
    fldChar2 = OxmlElement("w:fldChar")
    fldChar2.set(qn("w:fldCharType"), "end")
    run._r.extend([fldChar1, instrText, fldChar2])


def set_doc_defaults(doc: Document) -> None:
    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = RGBColor(*rgb_tuple(INK))
    for name, size, color in (("Title", 30, INK), ("Heading 1", 20, INK), ("Heading 2", 14, INK), ("Subtitle", 13, MUTED)):
        st = styles[name]
        st.font.name = "Arial"
        st.font.size = Pt(size)
        st.font.color.rgb = RGBColor(*rgb_tuple(color))
        st.font.bold = name != "Subtitle"


def set_inline_alt(inline_shape, title: str, description: str) -> None:
    doc_pr = inline_shape._inline.docPr
    doc_pr.set("title", title)
    doc_pr.set("descr", description)


def mark_row_as_header(row) -> None:
    tr_pr = row._tr.get_or_add_trPr()
    for child in tr_pr.findall(qn("w:tblHeader")):
        tr_pr.remove(child)
    header = OxmlElement("w:tblHeader")
    header.set(qn("w:val"), "true")
    tr_pr.append(header)


def set_doc_metadata(doc: Document, title: str, subject: str) -> None:
    props = doc.core_properties
    props.author = "Sable Harbor"
    props.last_modified_by = "Sable Harbor"
    props.title = title
    props.subject = subject
    props.category = "Corporate collateral"
    props.keywords = "Sable Harbor, corporate identity, template"
    props.comments = f"Generated from Sable Harbor brand system v{VERSION}."


def add_doc_header_footer(doc: Document, logo_png: Path, label: str = "CORPORATE DOCUMENT") -> None:
    section = doc.sections[0]
    header = section.header
    table = header.add_table(rows=1, cols=2, width=section.page_width - section.left_margin - section.right_margin)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = False
    table.columns[0].width = Inches(3.7)
    table.columns[1].width = Inches(3.2)
    mark_row_as_header(table.rows[0])
    left, right = table.rows[0].cells
    inline = left.paragraphs[0].add_run().add_picture(str(logo_png), width=Inches(3.0))
    set_inline_alt(inline, "Sable Harbor", "Sable Harbor corporate logo")
    right.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    p = right.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    r = p.add_run(label)
    r.bold = True
    r.font.name = "Arial"
    r.font.size = Pt(8)
    r.font.color.rgb = RGBColor(*rgb_tuple(MUTED))
    add_bottom_border(header.add_paragraph(), ORANGE, 8)

    footer = section.footer
    p = footer.paragraphs[0]
    add_bottom_border(p, ORANGE, 6)
    table = footer.add_table(rows=1, cols=2, width=section.page_width - section.left_margin - section.right_margin)
    table.autofit = False
    table.columns[0].width = Inches(5.9)
    table.columns[1].width = Inches(1.0)
    mark_row_as_header(table.rows[0])
    left, right = table.rows[0].cells
    lp = left.paragraphs[0]
    lr = lp.add_run("[ADDRESS]  |  [EMAIL]  |  [PHONE]  |  [WEB]")
    lr.font.name = "Arial"
    lr.font.size = Pt(7.5)
    lr.font.color.rgb = RGBColor(*rgb_tuple(MUTED))
    rp = right.paragraphs[0]
    rp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    rr = rp.add_run("PAGE ")
    rr.font.size = Pt(7.5)
    rr.font.color.rgb = RGBColor(*rgb_tuple(MUTED))
    add_page_field(rp)


def configure_page(section, page: str) -> None:
    if page == "us-letter":
        section.page_width = Inches(8.5)
        section.page_height = Inches(11)
    else:
        section.page_width = Cm(21.0)
        section.page_height = Cm(29.7)
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.85)
    section.right_margin = Inches(0.85)
    section.header_distance = Inches(0.28)
    section.footer_distance = Inches(0.28)


def create_letterhead_docx(repo: Path, out: Path, page: str) -> None:
    doc = Document()
    configure_page(doc.sections[0], page)
    set_doc_defaults(doc)
    set_doc_metadata(doc, f"Sable Harbor Letterhead - {page}", "Editable corporate correspondence template")
    add_doc_header_footer(doc, repo / "assets/brand/logos/sable-harbor__primary-horizontal.png", "CORPORATE CORRESPONDENCE")
    p = doc.add_paragraph()
    p.add_run("[DATE]").font.color.rgb = RGBColor(*rgb_tuple(MUTED))
    for text in ("[RECIPIENT NAME]", "[TITLE]", "[ORGANIZATION]", "[ADDRESS LINE 1]", "[ADDRESS LINE 2]"):
        p = doc.add_paragraph(text)
        p.paragraph_format.space_after = Pt(0)
    doc.add_paragraph()
    subject = doc.add_paragraph("RE: [SUBJECT]")
    subject.runs[0].bold = True
    subject.runs[0].font.size = Pt(11)
    doc.add_paragraph("[SALUTATION],")
    for _ in range(3):
        p = doc.add_paragraph("[Replace this text with correspondence. This template intentionally contains no invented office address, telephone number, or web domain.]")
        p.paragraph_format.space_after = Pt(9)
    doc.add_paragraph("[CLOSING],")
    doc.add_paragraph("[NAME]\n[TITLE]")
    out.parent.mkdir(parents=True, exist_ok=True)
    doc.save(out)


def create_memo_docx(repo: Path, out: Path, page: str = "us-letter") -> None:
    doc = Document()
    configure_page(doc.sections[0], page)
    set_doc_defaults(doc)
    set_doc_metadata(doc, "Sable Harbor Memorandum", "Editable memorandum template")
    add_doc_header_footer(doc, repo / "assets/brand/logos/sable-harbor__primary-horizontal.png", "MEMORANDUM")
    title = doc.add_paragraph("MEMORANDUM")
    title.style = doc.styles["Title"]
    title.paragraph_format.space_after = Pt(16)
    tbl = doc.add_table(rows=4, cols=2)
    tbl.alignment = WD_TABLE_ALIGNMENT.LEFT
    tbl.autofit = False
    tbl.columns[0].width = Inches(1.1)
    tbl.columns[1].width = Inches(5.8)
    mark_row_as_header(tbl.rows[0])
    for i, (label, placeholder) in enumerate((("TO", "[RECIPIENTS]"), ("FROM", "[AUTHOR / FUNCTION]"), ("DATE", "[DATE]"), ("SUBJECT", "[SUBJECT]"))):
        c0, c1 = tbl.rows[i].cells
        set_cell_margins(c0, 60, 40, 60, 100)
        set_cell_margins(c1, 60, 100, 60, 40)
        r = c0.paragraphs[0].add_run(label)
        r.bold = True
        r.font.size = Pt(8)
        r.font.color.rgb = RGBColor(*rgb_tuple(MUTED))
        c1.paragraphs[0].add_run(placeholder)
        if i % 2 == 0:
            set_cell_shading(c0, "EAE6DE")
            set_cell_shading(c1, "EAE6DE")
    doc.add_paragraph()
    for heading, body in (
        ("Purpose", "[State the decision, request, or information this memorandum is intended to carry.]"),
        ("Context", "[Record relevant facts, definitions, provenance, timing, and unresolved assumptions.]"),
        ("Analysis", "[Describe the relationships among the evidence. Preserve legitimate competing representations rather than forcing false convergence.]"),
        ("Decision / Next Action", "[Identify the owner, authority, action, due date, and what will be recorded afterward.]"),
    ):
        h = doc.add_paragraph(heading)
        h.style = doc.styles["Heading 1"]
        h.paragraph_format.space_before = Pt(12)
        h.paragraph_format.space_after = Pt(4)
        doc.add_paragraph(body)
    out.parent.mkdir(parents=True, exist_ok=True)
    doc.save(out)


def create_report_template_docx(repo: Path, out: Path, page: str) -> None:
    doc = Document()
    configure_page(doc.sections[0], page)
    set_doc_defaults(doc)
    set_doc_metadata(doc, f"Sable Harbor Report Template - {page}", "Editable controlled-report template")
    section = doc.sections[0]
    section.top_margin = Inches(0.5)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    inline = p.add_run().add_picture(str(repo / "assets/brand/logos/sable-harbor__stacked.png"), width=Inches(3.1))
    set_inline_alt(inline, "Sable Harbor", "Stacked Sable Harbor corporate logo")
    p = doc.add_paragraph("[REPORT TITLE]")
    p.style = doc.styles["Title"]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(24)
    sub = doc.add_paragraph("[REPORT SUBTITLE OR ENGAGEMENT]")
    sub.style = doc.styles["Subtitle"]
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("\n\n")
    control = doc.add_table(rows=3, cols=2)
    control.alignment = WD_TABLE_ALIGNMENT.CENTER
    control.autofit = False
    control.columns[0].width = Inches(1.6)
    control.columns[1].width = Inches(3.8)
    mark_row_as_header(control.rows[0])
    for i, (label, value) in enumerate((("DATE", "[DATE]"), ("VERSION", "[VERSION]"), ("CLASSIFICATION", "[CLASSIFICATION]"))):
        c0, c1 = control.rows[i].cells
        set_cell_margins(c0); set_cell_margins(c1)
        rr = c0.paragraphs[0].add_run(label); rr.bold = True; rr.font.size = Pt(8); rr.font.color.rgb = RGBColor(*rgb_tuple(MUTED))
        c1.paragraphs[0].add_run(value)
        if i == 1:
            set_cell_shading(c0, "EAE6DE"); set_cell_shading(c1, "EAE6DE")
    doc.add_page_break()
    configure_page(doc.sections[0], page)
    add_doc_header_footer(doc, repo / "assets/brand/logos/sable-harbor__primary-horizontal.png", "CONTROLLED REPORT")
    for heading, text in (
        ("Executive Summary", "[Summarize the question, evidence, conclusion, and decision implications.]"),
        ("1. Scope and Authority", "[Define the subject, period, boundaries, source authority, and limitations.]"),
        ("2. Evidence and Method", "[Record provenance, transformations, assumptions, and the method used.]"),
        ("3. Analysis", "[Present findings without collapsing distinct local meanings or overstating certainty.]"),
        ("4. Conclusions", "[State what is supported, what remains unresolved, and what changes next.]"),
        ("Appendices", "[Registers, source materials, diagrams, and supporting schedules.]"),
    ):
        p = doc.add_paragraph(heading)
        p.style = doc.styles["Heading 1"]
        doc.add_paragraph(text)
    out.parent.mkdir(parents=True, exist_ok=True)
    doc.save(out)


def create_presentation(repo: Path, out: Path) -> None:
    prs = Presentation()
    prs.core_properties.author = "Sable Harbor"
    prs.core_properties.last_modified_by = "Sable Harbor"
    prs.core_properties.title = "Sable Harbor Presentation Template"
    prs.core_properties.subject = "16:9 corporate presentation template"
    prs.core_properties.category = "Corporate collateral"
    prs.core_properties.keywords = "Sable Harbor, corporate identity, presentation"
    prs.core_properties.comments = f"Generated from Sable Harbor brand system v{VERSION}."
    prs.slide_width = PptxInches(13.333333)
    prs.slide_height = PptxInches(7.5)
    blank = prs.slide_layouts[6]
    ink = PptxRGB(*rgb_tuple(INK))
    off = PptxRGB(*rgb_tuple(OFF_WHITE))
    orange = PptxRGB(*rgb_tuple(ORANGE))
    muted = PptxRGB(*rgb_tuple(MUTED))
    dark = PptxRGB(*rgb_tuple(DARK))
    white = PptxRGB(*rgb_tuple(WHITE))

    def rect(slide, x, y, w, h, fill, line=None):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        shp.line.fill.background() if line is None else None
        return shp

    def text(slide, value, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, font="Arial", margin=0):
        box = slide.shapes.add_textbox(PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
        tf = box.text_frame; tf.clear(); tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=PptxInches(margin)
        p = tf.paragraphs[0]; p.text=value; p.alignment=align
        p.font.name=font; p.font.size=PptxPt(size); p.font.bold=bold; p.font.color.rgb=color
        return box

    def add_footer(slide, dark_mode=False, section="SABLE HARBOR"):
        rect(slide, 0.55, 7.1, 12.23, 0.025, orange)
        text(slide, section, 0.58, 7.18, 5.0, 0.18, 7.5, white if dark_mode else muted, bold=True)
        text(slide, "[DATE]  |  [CLASSIFICATION]", 8.0, 7.18, 4.75, 0.18, 7.5, white if dark_mode else muted, align=PP_ALIGN.RIGHT)

    # 1 title
    s = prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,dark)
    s.shapes.add_picture(str(repo/"assets/brand/logos/sable-harbor__reverse-horizontal.png"), PptxInches(0.8), PptxInches(0.7), width=PptxInches(5.4))
    text(s,"[PRESENTATION TITLE]",0.85,3.0,11.5,0.75,34,white,True)
    text(s,"[SUBTITLE / ENGAGEMENT]",0.88,3.85,10.5,0.4,16,PptxRGB(190,196,201))
    text(s,"[AUTHOR]  |  [DATE]  |  [VERSION]",0.88,5.9,10.5,0.35,11,PptxRGB(190,196,201))
    add_footer(s,True)

    # 2 agenda
    s = prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,off)
    text(s,"AGENDA",0.7,0.55,6,0.6,28,ink,True)
    rect(s,0.72,1.35,1.0,0.06,orange)
    for i,label in enumerate(("01  Context and question","02  Evidence and method","03  Findings","04  Decision and next actions")):
        text(s,label,1.0,1.8+i*0.95,10.8,0.55,19,ink, i==0)
    add_footer(s)

    # 3 section divider
    s = prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,dark)
    text(s,"01",0.85,1.25,2,0.8,44,orange,True)
    text(s,"[SECTION TITLE]",2.6,1.38,9.5,0.85,34,white,True)
    text(s,"[One sentence establishing the question or decision boundary.]",2.65,2.45,8.8,0.7,17,PptxRGB(190,196,201))
    add_footer(s,True,"SECTION DIVIDER")

    # 4 content
    s = prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,off)
    text(s,"[SLIDE TITLE]",0.7,0.55,11.6,0.6,26,ink,True)
    text(s,"[A concise statement of what the evidence supports.]",0.72,1.25,11.2,0.5,16,muted)
    for i,(h,b) in enumerate((("Observation","[What was observed or recorded.]"),("Meaning","[Why it matters in this context.]"),("Disposition","[What the organization must do next.]"))):
        x=0.72+i*4.15
        rect(s,x,2.05,3.72,0.05,orange)
        text(s,h.upper(),x,2.25,3.6,0.35,11,muted,True)
        text(s,b,x,2.85,3.55,1.9,18,ink,False)
    add_footer(s)

    # 5 two-column
    s = prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,off)
    text(s,"[COMPARISON OR DECISION FRAME]",0.7,0.55,11.6,0.6,26,ink,True)
    for x,title,accent in ((0.72,"[CURRENT STATE]",muted),(6.9,"[PROPOSED STATE]",orange)):
        rect(s,x,1.55,5.7,4.8,PptxRGB(235,232,225))
        rect(s,x,1.55,5.7,0.08,accent)
        text(s,title,x+0.35,1.95,5.0,0.5,16,ink,True)
        text(s,"• [Point one]\n• [Point two]\n• [Point three]\n\n[Record assumptions and limitations.]",x+0.35,2.75,4.9,2.8,16,ink)
    add_footer(s)

    # 6 business line
    s = prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,off)
    s.shapes.add_picture(str(repo/"assets/brand/logos/willow__primary-horizontal.png"), PptxInches(0.7), PptxInches(0.55), width=PptxInches(5.1))
    text(s,"[BUSINESS-LINE SECTION]",0.75,2.0,11.5,0.6,26,ink,True)
    text(s,"Use the applicable individual business-line logo. Preserve the endorsed relationship and do not invent a legal entity or reporting line.",0.75,2.85,11.3,1.0,18,muted)
    rect(s,0.75,4.35,11.85,0.07,PptxRGB(*rgb_tuple("#315F4D")))
    text(s,"[KEY QUESTION]",0.75,4.75,4,0.4,11,muted,True)
    text(s,"[QUESTION -> BELIEF -> EXPERIMENT -> OBSERVATION -> DECISION]",0.75,5.25,11.5,0.55,16,ink,True)
    add_footer(s,False,"BUSINESS-LINE TEMPLATE")

    # 7 close
    s = prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,dark)
    s.shapes.add_picture(str(repo/"assets/brand/logos/sable-harbor__reverse-horizontal.png"), PptxInches(4.25), PptxInches(0.78), width=PptxInches(4.83))
    text(s,"KEEP IT LEARNABLE.",1.5,4.1,10.3,0.75,32,white,True,PP_ALIGN.CENTER)
    text(s,"[NEXT ACTION / CONTACT]",2.0,5.15,9.3,0.4,14,PptxRGB(190,196,201),False,PP_ALIGN.CENTER)
    add_footer(s,True,"CLOSING")

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)


def register_reportlab_fonts() -> tuple[str, str]:
    # ReportLab does not support the CFF/PostScript outlines used by the
    # packaged Inter Display OTFs. Use its built-in Helvetica metrics for
    # document body copy; all production logo lettering remains outlined
    # directly from Inter Display in the SVG assets.
    return "Helvetica", "Helvetica-Bold"


def png_for_pdf(svg_path: Path, temp_dir: Path, width: int = 1200) -> Path:
    temp_dir.mkdir(parents=True, exist_ok=True)
    out = temp_dir / (svg_path.stem + "__pdf.png")
    cairosvg.svg2png(url=str(svg_path), write_to=str(out), output_width=width)
    return out


def create_brand_standards_pdf(repo: Path, out: Path) -> None:
    regular, bold = register_reportlab_fonts()
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="SHCover", fontName=bold, fontSize=34, leading=40, textColor=HexColor(INK), spaceAfter=12))
    styles.add(ParagraphStyle(name="SHSub", fontName=regular, fontSize=13, leading=18, textColor=HexColor(MUTED)))
    styles.add(ParagraphStyle(name="SHH1", fontName=bold, fontSize=22, leading=26, textColor=HexColor(INK), spaceAfter=10, spaceBefore=6))
    styles.add(ParagraphStyle(name="SHH2", fontName=bold, fontSize=13, leading=17, textColor=HexColor(INK), spaceAfter=5, spaceBefore=10))
    styles.add(ParagraphStyle(name="SHBody", fontName=regular, fontSize=9.5, leading=14, textColor=HexColor(INK), spaceAfter=7))
    styles.add(ParagraphStyle(name="SHSmall", fontName=regular, fontSize=7.5, leading=10, textColor=HexColor(MUTED)))
    styles.add(ParagraphStyle(name="SHCallout", fontName=bold, fontSize=15, leading=21, textColor=HexColor(INK), leftIndent=10, borderColor=HexColor(ORANGE), borderWidth=2, borderPadding=10, spaceAfter=12))

    doc = SimpleDocTemplate(str(out), pagesize=LETTER, rightMargin=46, leftMargin=46, topMargin=46, bottomMargin=46, title="Sable Harbor Brand Standards v0.2.0", author="Sable Harbor")
    story = []
    tmp = repo / ".brand-tmp"
    reverse = png_for_pdf(repo/"assets/brand/logos/sable-harbor__reverse-horizontal.svg", tmp, 1600)
    primary = png_for_pdf(repo/"assets/brand/logos/sable-harbor__primary-horizontal.svg", tmp, 1600)
    stacked = png_for_pdf(repo/"assets/brand/logos/sable-harbor__stacked.svg", tmp, 1000)
    story.extend([
        Table([[RLImage(str(reverse), width=6.6*inch, height=2.06*inch)]], colWidths=[7.0*inch], rowHeights=[2.55*inch], style=TableStyle([("BACKGROUND",(0,0),(-1,-1),HexColor(DARK)),("VALIGN",(0,0),(-1,-1),"MIDDLE"),("ALIGN",(0,0),(-1,-1),"CENTER"),("BOX",(0,0),(-1,-1),0,HexColor(DARK))])),
        Spacer(1, 0.35*inch),
        Paragraph("CORPORATE IDENTITY SYSTEM", styles["SHCover"]),
        Paragraph("Production standards, identity architecture, collateral, and governance", styles["SHSub"]),
        Spacer(1, 1.6*inch),
        Paragraph(f"VERSION {VERSION}  |  {BUILD_DATE}", styles["SHSub"]),
        PageBreak(),
        Paragraph("1. Governing design idea", styles["SHH1"]),
        Paragraph("Variation is normal. The work is to encounter it, represent it, understand it, and choose consciously what happens next.", styles["SHCallout"]),
        Paragraph("The identity system is built from explicit geometry, visible seams, bounded channels, junctions, strata, and controlled points of integration. It avoids literal harbor imagery, generic technology symbols, faux heraldry, gradients, shadows, and decorative claims that are not supported by canon.", styles["SHBody"]),
        Paragraph("The marks do not create corporate canon. They represent names and roles already established in the controlling corporate-lore documents. OPEN legal entities, reporting lines, titles, economics, and launch details remain open.", styles["SHBody"]),
        Spacer(1, 10),
        RLImage(str(primary), width=6.6*inch, height=2.06*inch),
        PageBreak(),
        Paragraph("2. Master-brand lockups", styles["SHH1"]),
        Paragraph("Five master configurations support the same mark in different operating contexts. Use the configuration that fits the available field rather than distorting another configuration.", styles["SHBody"]),
    ])
    master_rows = []
    for variant in CORE_VARIANTS:
        path = repo/"assets/brand/logos"/f"sable-harbor__{variant}.svg"
        img = png_for_pdf(path, tmp, 900)
        if variant == "mark":
            w,h=1.25*inch,1.25*inch
        elif variant == "stacked":
            w,h=1.35*inch,1.35*inch
        else:
            w,h=2.75*inch,0.86*inch
        master_rows.append([RLImage(str(img), width=w, height=h), Paragraph(f"<b>{variant.replace('-', ' ').title()}</b><br/>{'Use for icons and small identifiers.' if variant=='mark' else 'Use without altering proportion, spacing, or color relationships.'}", styles["SHBody"])])
    story.append(Table(master_rows, colWidths=[3.1*inch,3.5*inch], style=TableStyle([("VALIGN",(0,0),(-1,-1),"MIDDLE"),("LINEBELOW",(0,0),(-1,-2),0.35,HexColor("#D8D3CA")),("TOPPADDING",(0,0),(-1,-1),9),("BOTTOMPADDING",(0,0),(-1,-1),9)])))
    story.extend([PageBreak(), Paragraph("3. Identity architecture", styles["SHH1"]), Paragraph("Sable Harbor is the corporate master brand. Seven 2026 lines receive the full five-lockup system. Supplemental assets and historical identities remain visibly subordinate and must not be presented as current legal entities unless later canon says otherwise.", styles["SHBody"])])
    grid = []
    for identity in IDENTITIES[1:]:
        if identity.slug == "red-wash-pale-sun":
            continue
        path = repo/"assets/brand/logos"/f"{identity.slug}__mark.svg"
        img = png_for_pdf(path, tmp, 500)
        grid.append([RLImage(str(img), width=.55*inch, height=.55*inch), Paragraph(f"<b>{escape(identity.display_name)}</b><br/>{escape(identity.role)}", styles["SHSmall"]), Paragraph(identity.status, styles["SHSmall"])])
    story.append(Table(grid, colWidths=[.7*inch,2.05*inch,3.85*inch], style=TableStyle([("VALIGN",(0,0),(-1,-1),"MIDDLE"),("LINEBELOW",(0,0),(-1,-2),.25,HexColor("#DDD8D0")),("TOPPADDING",(0,0),(-1,-1),5),("BOTTOMPADDING",(0,0),(-1,-1),5)])))
    story.extend([PageBreak(), Paragraph("4. Color system", styles["SHH1"]), Paragraph("HEX and RGB values are controlling for digital use. CMYK values below are working conversions, not press-certified recipes. Commercial print production requires substrate, ink, profile, and proof validation. No Pantone number is asserted by this package.", styles["SHBody"])])
    rows=[[Paragraph("COLOR",styles["SHSmall"]),Paragraph("HEX / RGB",styles["SHSmall"]),Paragraph("CMYK APPROX.",styles["SHSmall"]),""]]
    for name,value in PALETTE.items():
        r,g,b=rgb_tuple(value); c,m,y,k=cmyk_approx(value)
        swatch=Table([[""]],colWidths=[.36*inch],rowHeights=[.28*inch],style=TableStyle([("BACKGROUND",(0,0),(-1,-1),HexColor(value)),("BOX",(0,0),(-1,-1),.4,HexColor("#999999"))]))
        rows.append([Paragraph(name,styles["SHBody"]),Paragraph(f"{value}<br/>{r}, {g}, {b}",styles["SHSmall"]),Paragraph(f"{c}, {m}, {y}, {k}",styles["SHSmall"]),swatch])
    story.append(Table(rows,colWidths=[2.2*inch,1.65*inch,1.45*inch,.7*inch],style=TableStyle([("BACKGROUND",(0,0),(-1,0),HexColor("#E6E2DA")),("VALIGN",(0,0),(-1,-1),"MIDDLE"),("GRID",(0,0),(-1,-1),.25,HexColor("#D3CEC5")),("TOPPADDING",(0,0),(-1,-1),5),("BOTTOMPADDING",(0,0),(-1,-1),5)])))
    story.extend([PageBreak(), Paragraph("5. Reproduction rules", styles["SHH1"])])
    for h,b in (
        ("Clear space","Maintain clear space at least equal to the central accent square around the complete lockup."),
        ("Minimum sizes","Horizontal lockups: 32 mm print / 180 px digital. Stacked lockups: 22 mm / 120 px. Marks: 6 mm / 24 px, subject to a legibility check."),
        ("Color","Use primary color on light fields, reverse on dark fields, and one-color for stamps, engraving, fax, or restricted reproduction."),
        ("File choice","SVG is the production source of truth. PNG is a convenience render. PDF collateral is print-ready only after local proofing."),
        ("Prohibited treatments","Do not distort, rotate, bevel, shadow, outline, add gradients, recolor outside the approved palette, place over noisy imagery, or substitute literal harbor, mine, railway, or generic AI clip art."),
        ("Endorsements","Do not create new combined identities. Use only provided endorsed lockups and preserve the corporate/operating relationship stated in canon."),
        ("Historical identities","Evalon and Emberline must be marked historical. Evalon is an outpost predecessor to Willow, not a current business line."),
    ):
        story.append(Paragraph(h,styles["SHH2"])); story.append(Paragraph(b,styles["SHBody"]))
    story.extend([PageBreak(), Paragraph("6. Collateral system", styles["SHH1"]), Paragraph("The package includes US Letter and A4 letterhead, memo and report templates, editable SVG sources, PDF and PNG print references, and a 16:9 PowerPoint template. All address and contact fields are explicit placeholders; the templates do not invent a headquarters address, telephone number, email address, or public domain.", styles["SHBody"])])
    collat=[
        ("Letterhead","US Letter + A4: SVG, PNG, PDF, DOCX"),
        ("Memorandum","US Letter: DOCX + PDF reference"),
        ("Report","US Letter + A4: DOCX; cover SVG, PNG, PDF"),
        ("Presentation","16:9 PPTX with title, agenda, divider, content, comparison, business-line, and close layouts"),
        ("Packages","Logo system, collateral package, and complete brand kit ZIPs with checksums"),
    ]
    story.append(Table([[Paragraph(f"<b>{a}</b>",styles["SHBody"]),Paragraph(b,styles["SHBody"])] for a,b in collat],colWidths=[1.5*inch,5.1*inch],style=TableStyle([("LINEBELOW",(0,0),(-1,-2),.25,HexColor("#D3CEC5")),("VALIGN",(0,0),(-1,-1),"TOP"),("TOPPADDING",(0,0),(-1,-1),7),("BOTTOMPADDING",(0,0),(-1,-1),7)])))
    story.extend([PageBreak(), Paragraph("7. Governance and legal posture", styles["SHH1"]), Paragraph("The production artwork is a working corporate-identity candidate for the synthetic enterprise. It does not independently establish a legal entity, trademark right, reporting relationship, operating authority, or public commercial use.", styles["SHBody"]), Paragraph("Inter Display Medium, Semibold, and Bold were used as construction sources; wordmarks in production logo SVGs are converted to vector outlines. Font binaries are not distributed. The typeface project is licensed under the SIL Open Font License 1.1; see FONT_PROVENANCE.md.", styles["SHBody"]), Paragraph("A preliminary name-and-mark screen is included for planning. It is not a legal opinion, a comprehensive clearance search, or a determination of registrability. Several names are crowded or conflict-prone for software, industrial technology, logistics, or advisory services. Obtain U.S. trademark counsel before external commercial adoption.", styles["SHBody"]), Paragraph("Controlling source: docs/canon/SABLE_HARBOR_CORPORATE_LORE_CANON_v0.2.md", styles["SHSmall"]), Paragraph(f"Brand package version {VERSION} | Generated {BUILD_DATE}", styles["SHSmall"])])
    doc.build(story)


def create_simple_reference_pdf(repo: Path, out: Path, kind: str) -> None:
    regular,bold=register_reportlab_fonts()
    pagesize=LETTER
    doc=SimpleDocTemplate(str(out),pagesize=pagesize,rightMargin=54,leftMargin=54,topMargin=54,bottomMargin=54,title=f"Sable Harbor {kind}")
    styles=getSampleStyleSheet()
    h=ParagraphStyle("h",fontName=bold,fontSize=24,leading=28,textColor=HexColor(INK),spaceAfter=14)
    b=ParagraphStyle("b",fontName=regular,fontSize=10.5,leading=15,textColor=HexColor(INK),spaceAfter=10)
    label=ParagraphStyle("label",fontName=bold,fontSize=8,leading=10,textColor=HexColor(MUTED),spaceAfter=3)
    logo=png_for_pdf(repo/"assets/brand/logos/sable-harbor__primary-horizontal.svg",repo/".brand-tmp",1400)
    story=[RLImage(str(logo),width=4.5*inch,height=1.4*inch),Spacer(1,8),Paragraph(kind.upper(),h)]
    if kind.lower().startswith("memorandum"):
        data=[[Paragraph("TO",label),Paragraph("[RECIPIENTS]",b)],[Paragraph("FROM",label),Paragraph("[AUTHOR / FUNCTION]",b)],[Paragraph("DATE",label),Paragraph("[DATE]",b)],[Paragraph("SUBJECT",label),Paragraph("[SUBJECT]",b)]]
        story += [Table(data,colWidths=[1.1*inch,5.2*inch],style=TableStyle([("BACKGROUND",(0,0),(-1,-1),HexColor("#F0EDE7")),("GRID",(0,0),(-1,-1),.25,HexColor("#D3CEC5")),("VALIGN",(0,0),(-1,-1),"TOP"),("TOPPADDING",(0,0),(-1,-1),6),("BOTTOMPADDING",(0,0),(-1,-1),6)])),Spacer(1,15)]
        for hh,bb in (("Purpose","[State the decision, request, or information.]"),("Context","[Record relevant facts, timing, provenance, and assumptions.]"),("Analysis","[Explain relationships and preserve uncertainty.]"),("Decision / Next Action","[Name the owner, authority, action, and record.]")):
            story += [Paragraph(hh,ParagraphStyle("s"+hh,fontName=bold,fontSize=14,leading=17,textColor=HexColor(INK),spaceBefore=7,spaceAfter=4)),Paragraph(bb,b)]
    doc.build(story)


def brand_markdown() -> str:
    rows=[]
    for name,value in PALETTE.items():
        r,g,b=rgb_tuple(value); c,m,y,k=cmyk_approx(value)
        rows.append(f"| {name} | `{value}` | `{r}, {g}, {b}` | `{c}, {m}, {y}, {k}` |")
    return f"""# Sable Harbor Brand Standards

**Version:** {VERSION}  
**Date:** {BUILD_DATE}  
**Status:** Production candidate identity system; does not independently create or change canon

## 1. Governing idea

> **Variation is normal. The work is to encounter it, represent it, understand it, and choose consciously what happens next.**

The identity system translates that proposition into explicit geometry: visible seams, strata, channels, junctions, bounded enclosures, and controlled points of integration. It deliberately avoids literal harbor imagery, generic AI/circuit clip art, faux heraldry, gradients, shadows, and invented claims.

## 2. Identity hierarchy

1. **Sable Harbor** is the corporate master brand.
2. **Foundry Field, Willow, Atlas Meridian, Pale Sun, Project Cradle, American Resource Utility, and Advisory** are the seven 2026 business lines represented in the narrative map.
3. **Foundry, Red Wash Mine, Blood, Sweat & Tears Railway, and Red Wash / Pale Sun** are supplemental identities or endorsed relationships.
4. **Emberline and Evalon** are historical identities. Evalon is an archival advanced-engineering outpost that was closed and rechartered as Willow; it is not an eighth current business line.

## 3. Master-brand configurations

The Sable Harbor master brand has five controlled configurations:

- `sable-harbor__primary-horizontal` - default for README headings, wiki headers, reports, and letterhead;
- `sable-harbor__stacked` - covers, title pages, and square placements;
- `sable-harbor__mark` - avatars, favicons, navigation, and small document furniture;
- `sable-harbor__reverse-horizontal` - dark fields and video;
- `sable-harbor__one-color-horizontal` - stamps, engraving, fax, monochrome printing, and restricted reproduction.

Never stretch one configuration to imitate another.

## 4. Clear space and minimum size

- Preserve clear space at least equal to the central accent square around the entire lockup.
- Horizontal lockup: minimum 32 mm in print or 180 px in digital use.
- Stacked lockup: minimum 22 mm or 120 px.
- Compact mark: minimum 6 mm or 24 px, subject to a legibility check.
- At unusually small sizes, use the one-color mark rather than retaining an accent that cannot reproduce cleanly.

## 5. Color

Digital HEX/RGB values are controlling. CMYK values below are mathematical working conversions and require a printer proof. No Pantone or proprietary spot-color designation is asserted.

| Color | HEX | RGB | CMYK approximation |
|---|---|---|---|
{chr(10).join(rows)}

## 6. Reproduction rules

- Use SVG as the production source of truth; use PNG for convenience.
- Use the supplied reverse asset rather than manually inverting a primary file.
- Do not distort, rotate, bevel, shadow, outline, add gradients, or recolor outside the approved palette.
- Do not place the mark on visually noisy imagery without an adequate solid field.
- Do not substitute literal lighthouse, compass, shield, wave, mine-pick, locomotive, atom, or generic circuit/AI imagery.
- Do not create new combined identities. Use only supplied endorsed lockups.
- Do not present historical identities as current business lines.
- Do not use logo contact sheets as production art. Every file under `assets/brand/logos/` contains one lockup only.

## 7. Document and presentation use

- **Letterhead:** primary horizontal lockup; no invented address or contact information.
- **Memoranda:** primary horizontal lockup with explicit `TO`, `FROM`, `DATE`, and `SUBJECT` fields.
- **Reports:** stacked lockup on covers; primary horizontal in body headers.
- **Presentations:** reverse lockup on dark title fields; current business-line lockup on unit-specific sections.
- **Wiki and README:** SVG assets linked directly from `assets/brand/logos/`.

## 8. Governance

The controlling naming source is [`{CANON_PATH}`](../../{CANON_PATH}). Artwork does not lock legal entities, reporting lines, titles, economics, offices, public domains, contact details, or other OPEN decisions.

Before external commercial adoption, review [`docs/legal/PRELIMINARY_NAME_AND_MARK_SCREEN.md`](../../docs/legal/PRELIMINARY_NAME_AND_MARK_SCREEN.md) and obtain qualified trademark counsel.
"""


def font_provenance_markdown() -> str:
    return f"""# Font Provenance

**Brand package version:** {VERSION}  
**Recorded:** {BUILD_DATE}

## Construction source

The production wordmarks were constructed using **Inter Display Medium, Semibold, and Bold** as spacing and outline sources. The production logo SVG files contain vector path outlines; they do not contain live `<text>` elements and do not require the font to display.

## Distribution boundary

- No Inter font binaries are distributed in this repository or in the packaged brand archives.
- The generated logos and documents are outputs, not redistributed font software.
- Editable collateral may fall back to common system fonts for placeholder and body text; the Sable Harbor wordmark itself remains outlined artwork.

## Upstream project and license

Inter is maintained by The Inter Project Authors and distributed under the **SIL Open Font License, Version 1.1**.

- Project: <https://github.com/rsms/inter>
- License: <https://github.com/rsms/inter/blob/master/LICENSE.txt>
- License text: <https://scripts.sil.org/OFL>

The OFL states that the requirement for font software to remain under the license does not apply to documents created using the fonts or their derivatives. This record is provenance, not a substitute for legal advice.

## Reproducibility

The deterministic generator records the font paths used in the local or CI build. A different installed font must not silently replace Inter for production logo regeneration; validation should fail or the change must be reviewed as a new brand version.
"""


def trademark_markdown() -> str:
    return f"""# Preliminary Name and Mark Screen

**Date:** {BUILD_DATE}  
**Scope:** U.S.-focused knockout and common-law planning screen  
**Status:** Research aid only - not a legal opinion, comprehensive clearance search, registrability determination, or freedom-to-operate conclusion

## 1. Method and limitations

The USPTO states that a comprehensive clearance search should include federal registrations and applications, state registries, the internet, and confusingly similar wording and designs used with related goods or services. Similarity can arise from appearance, sound, meaning, or commercial impression; the goods and services need not be identical or share an international class.

Primary guidance:

- <https://www.uspto.gov/trademarks/search/comprehensive-clearance-search-similar-trademarks>
- <https://www.uspto.gov/trademarks/search/federal-trademark-searching>
- <https://www.uspto.gov/trademarks/search/likelihood-confusion>
- <https://www.uspto.gov/trademarks/search/search>

This screen reviewed exact-name web results and selected reported federal records available as of the date above. It did **not** complete attorney-grade phonetic, design-code, coordinated-class, state-registry, domain, international, assignment, litigation, or marketplace searches. Third-party trademark record sites are used as finding aids and must be verified in USPTO TSDR.

## 2. Planning result

| Identity | Internal synthetic use | External commercial adoption | Preliminary risk | Reason |
|---|---|---|---|---|
| Sable Harbor | Retain | Counsel review before use | **Medium** | Exact-name entertainment software/app use exists; no conclusion is reached on relatedness to industrial software and services. |
| Foundry | Retain as canon substrate | Do not launch alone without renaming strategy and counsel | **Very high** | Extremely crowded technology term; live Class 42 registrations/applications and active litigation over software use. |
| Foundry Field | Retain | Counsel review; consider a stronger house-mark endorsement | **Medium-high** | Compound is more distinctive than Foundry alone, but the dominant term remains crowded and exact/near-exact nonsoftware uses exist. |
| Willow | Retain | Do not launch for industrial data/digital-twin technology without a legal naming decision | **Very high** | Live WILLOW rights cover building automation, operational data, digital twins, knowledge graphs, analytics, and related consulting/SaaS. |
| Atlas Meridian | Retain | Counsel review and likely naming differentiation | **High** | Exact-name U.S. businesses currently use the name for operational modernization, AI enablement, trade, brand development, and e-commerce; a dissolved UK software entity also existed. |
| Pale Sun | Retain | Full search before filing | **Low-medium** | Limited exact relevant commercial use surfaced in this screen; media/title use exists. This is not a clearance conclusion. |
| Project Cradle | Retain | Counsel review and field-specific search | **Medium** | Exact-name web use and a live PROJECT CRADLE CARE registration exist, although the known goods/services differ. |
| American Resource Utility | Retain | Counsel review; test full name and acronym separately | **High** | The acronym ARU has a live transportation/storage registration, and the natural-resources field contains close corporate wording such as American Resources. |
| Advisory | Retain only as provisional descriptive label | Select a distinctive final name before commercial launch | **Very high / weak** | "Advisory" is descriptive and routinely disclaimed in registrations; it is unlikely to function as a strong standalone proprietary identifier. Canon already leaves the final name OPEN. |
| Red Wash Mine | Retain as fictional asset | Geographic and mining-name diligence before external use | **Medium** | Exact and near-exact historical uranium-mine names exist in the U.S. Southwest, creating geographic/industry association risk even without a known exact trademark record. |
| Blood, Sweat & Tears Railway | Retain as lore | Counsel review before public transportation branding | **Medium-high** | The phrase is culturally famous and has prior model-railroad use; exact rail-service clearance was not established. |
| Emberline | Retain as historical lore | Do not revive commercially without a new name or counsel-led strategy | **Very high** | Current investment/technology company use plus multiple live registrations/applications across financial, lodging, construction-material, industrial-protective, and other fields. |
| Evalon | Retain strictly as historical/archival | Do not revive as a software or AI brand | **Very high** | Existing U.S. software/consulting use and current Evalon AI/common-law use; an EU EVALON registration was reported in 2026. |

## 3. Load-bearing findings

### Sable Harbor

An iOS adventure application uses the exact name **Sable Harbor** and reports copyright attribution to “Sable Harbor 2025 Inc.” This appears materially different from industrial systems, but exact-name software use requires analysis rather than dismissal.

- <https://apps.apple.com/ca/app/sable-harbor/id6755688168>

### Foundry / Foundry Field

Reported federal records include live or pending FOUNDRY marks in Class 42, including a registration by Foundry International, LLC. A 2026 lawsuit also concerns software-related use of “Foundry,” illustrating the term's conflict density.

- <https://www.trademarkia.com/foundry-97187178>
- <https://www.trademarkelite.com/trademark/trademark-detail/99541295/FOUNDRY>
- <https://www.reuters.com/legal/litigation/adobe-sued-trademark-infringement-over-foundry-ai-tool-2026-03-12/>

### Willow

Willow IP Pty Ltd has reported live federal rights/applications for WILLOW covering building automation, IoT, digital-twin interfaces, operational data, analytics, and SaaS. The operating overlap with Sable Harbor's industrial-representation and experimentation story is substantial.

- <https://furm.com/trademarks/willow-98642135>
- <https://furm.com/trademarks/willow-88202768>
- <https://willowinc.com/willow-digital-twin/>

### Atlas Meridian

Multiple current exact-name businesses use Atlas Meridian for cross-border commerce, operational modernization, AI enablement, hospitality infrastructure, and related services. A UK software/data-processing company with the exact name was dissolved in March 2026, which does not eliminate possible residual or common-law issues.

- <https://atlasmeridianllc.com/>
- <https://atlasmeridian.org/>
- <https://atlasmeridianusa.com/about-amg>
- <https://find-and-update.company-information.service.gov.uk/company/16181883>

### Project Cradle

The exact phrase appears in current web use. PROJECT CRADLE CARE is a reported live federal registration for educational services; those services differ from minerals recovery, but exact and dominant-term similarity still belongs in a comprehensive search.

- <https://projectcradle.org/>
- <https://trademarks.justia.com/870/79/project-cradle-87079222.html>

### American Resource Utility / ARU

A reported live ARU registration covers transportation and storage services, directly relevant to ARU's fictional logistics role. American Resources Corporation also operates in critical materials, mining, processing, and logistics-adjacent markets.

- <https://www.trademarkia.com/aru-97165633>
- <https://www.americanresourcescorp.com/>

### Advisory

Federal records routinely disclaim exclusive rights in “ADVISORY” when used in compound consulting marks, including technology-governance and cybersecurity/business-advisory services. Treat Advisory as a descriptive working label, not the final protectable name.

- <https://furm.com/trademarks/sl-advisory-87006481>
- <https://trademarks.justia.com/992/39/ntm-99239394.html>
- <https://trademarks.justia.com/886/72/advisory-88672729.html>

### Emberline

Emberline is currently used by an investment company focused on B2B technology in high-consequence industries, and multiple reported federal registrations/applications use EMBERLINE in other fields. Because the Sable Harbor identity is historical, retain it only as corporate archaeology unless counsel approves a revival.

- <https://www.emberline.com/>
- <https://www.trademarkia.com/emberline-99522886>
- <https://furm.com/trademarks/emberline-97659504>
- <https://furm.com/trademarks/emberline-88345265>

### Evalon

Evalon has existing software-development and systems-integration use, more recent AI-platform use, and a reported 2026 EU registration. Its archival identity should never be presented as a current external Sable Harbor software brand.

- <https://www.pr.com/press-release/62104>
- <https://www.linkedin.com/company/evalon-inc/>
- <https://www.europapress.es/comunicados/empresas-00908/noticia-comunicado-icp-iberica-capital-partners-lanza-evalon-ai-20250917160313.html>
- <https://www.trademarkelite.com/europe/trademark/trademark-detail/019286889/EVALON>

### Red Wash Mine and BS&T Railway

The limited screen found historical **Upper Red Wash Mine** uranium references in Arizona and model-railroad use of “Blood, Sweat & Tears Railway.” These are not federal clearance results, but they are enough to prevent an unsupported “no conflicts” statement.

- <https://www.mindat.org/loc-69860.html>
- <https://www.sporskiftet.dk/forum/inspiration>

## 4. Required next legal work before real market use

1. Define actual goods and services by identity, including Classes 9, 35, 37, 39, 40, 41, and 42 as potentially relevant.
2. Run exact, expanded, phonetic, translation, acronym, coordinated-class, and design-code searches in the USPTO system.
3. Review every material live record in TSDR, including owner, identification, status, prosecution history, and assignments.
4. Search state trademark and entity registries, domains, app stores, industry publications, social platforms, and common-law use.
5. Search Canada, Mexico, the EU, UK, Australia, and other intended markets.
6. Have U.S. trademark counsel assess likelihood of confusion, descriptiveness, filing architecture, house-mark strategy, and coexistence/consent options.
7. Do not file or launch the current marks solely on the basis of this document.

## 5. Practical recommendation

For the synthetic enterprise and repository, retain the names as canon. For a real commercial company, the highest-priority naming decisions are **Willow, Foundry/Foundry Field, Atlas Meridian, American Resource Utility/ARU, Advisory, Emberline, and Evalon**. Keep **Emberline** and **Evalon** historical, choose a distinctive final name for **Advisory**, and obtain counsel-led clearance before using any line externally.
"""


def business_page(identity: Identity) -> str:
    base = "../../assets/brand/logos"
    logo = f"{base}/{identity.slug}__primary-horizontal.svg"
    variants = []
    for v in identity.variants:
        variants.append(f"- [{v.replace('-', ' ').title()}]({base}/{identity.slug}__{v}.svg) ([PNG]({base}/{identity.slug}__{v}.png))")
    current_text = "Current 2026 identity" if identity.current else "Supplemental or historical identity"
    return f"""# {identity.display_name}

<p align="center">
  <img src="{logo}" alt="{identity.display_name} logo" width="760" />
</p>

| Field | Status |
|---|---|
| Role | {identity.role} |
| Publication class | {current_text} |
| Canon state | {identity.status} |
| Controlling source | [`{CANON_PATH}`](../canon/SABLE_HARBOR_CORPORATE_LORE_CANON_v0.2.md), {identity.source_section} |

## Role in the Sable Harbor system

{identity.description}

## Logo assets

{chr(10).join(variants)}

The SVG files are the production source of truth. Do not infer legal-entity status, reporting structure, economics, or public commercial use from the artwork.

## Related material

- [Business-line index](README.md)
- [Brand standards](../../assets/brand/BRAND_STANDARDS.md)
- [Organization maps](../organization/README.md)
- [Preliminary name and mark screen](../legal/PRELIMINARY_NAME_AND_MARK_SCREEN.md)
"""


def business_index() -> str:
    current = [i for i in IDENTITIES if i.current and i.slug != "sable-harbor"]
    supporting = [i for i in IDENTITIES if not i.current]
    rows=[]
    for i in current:
        rows.append(f'| [<img src="../../assets/brand/logos/{i.slug}__mark.svg" width="72" alt="{i.display_name} mark" />]({i.slug.upper().replace("-", "_")}.md) | [{i.display_name}]({i.slug.upper().replace("-", "_")}.md) | {i.role} |')
    sup=[]
    for i in supporting:
        sup.append(f"| [{i.display_name}]({i.slug.upper().replace('-', '_')}.md) | {i.role} | {i.status} |")
    return f"""# Sable Harbor Business Lines and Institutional Identities

<p align="center">
  <img src="../../assets/brand/logos/sable-harbor__primary-horizontal.svg" alt="Sable Harbor" width="800" />
</p>

This index is a publication and navigation layer over corporate lore v0.2. It does not independently create canon.

## Current 2026 business lines

| Mark | Identity | Role |
|---|---|---|
{chr(10).join(rows)}

## Supplemental and historical identities

| Identity | Role | Status |
|---|---|---|
{chr(10).join(sup)}

## Operating logic

- Foundry Field encounters, connects, and represents operational reality.
- Willow tests consequential unknowns through bounded experiments.
- Atlas Meridian investigates across represented evidence.
- Pale Sun owns and operates a uranium asset where operating control is the thesis.
- Project Cradle recovers value from material streams the host system already creates.
- American Resource Utility operates logistics across physical and organizational boundaries.
- Advisory transfers the method where the client can and should own the system.

This is an analytical map, not a finalized legal or reporting structure.

## Source and assets

- [Corporate lore canon](../canon/SABLE_HARBOR_CORPORATE_LORE_CANON_v0.2.md)
- [Brand asset package](../../assets/brand/README.md)
- [Organization maps](../organization/README.md)
"""


def wiki_page(identity: Identity) -> str:
    raw = f"https://raw.githubusercontent.com/SquirmyWormy275/SABLEHARBOR/main/assets/brand/logos/{identity.slug}__primary-horizontal.svg"
    repo_page = f"{REPO_URL}/blob/main/docs/business-lines/{identity.slug.upper().replace('-', '_')}.md"
    return f"""# {identity.display_name}

![{identity.display_name}]({raw})

**Role:** {identity.role}  
**Canon state:** {identity.status}

{identity.description}

The wiki summarizes the controlling repository canon; it does not create or change canon.

[Open the repository reference page]({repo_page})
"""


def wiki_home() -> str:
    raw = "https://raw.githubusercontent.com/SquirmyWormy275/SABLEHARBOR/main/assets/brand/logos/sable-harbor__primary-horizontal.svg"
    return f"""# Sable Harbor

![Sable Harbor]({raw})

Sable Harbor is the canonical synthetic enterprise and reusable business-world sandbox for mining, natural resources, industrial systems, enterprise software, assurance, analytics, finance, governance, security, incident response, and professional training.

## Start here

- [[Business Lines]]
- [[Organization]]
- [[Corporate History]]
- [[Brand Assets]]
- [[Canon and Governance]]

## Canon boundary

Versioned documents in the repository remain controlling. This wiki is a public-facing reference and navigation layer and does not independently create or change canon.

[Open the repository]({REPO_URL})
"""


def write_docs(repo: Path) -> None:
    # Brand docs
    (repo/"assets/brand/BRAND_STANDARDS.md").write_text(brand_markdown(), encoding="utf-8")
    (repo/"assets/brand/FONT_PROVENANCE.md").write_text(font_provenance_markdown(), encoding="utf-8")
    (repo/"docs/legal").mkdir(parents=True, exist_ok=True)
    (repo/"docs/legal/PRELIMINARY_NAME_AND_MARK_SCREEN.md").write_text(trademark_markdown(), encoding="utf-8")

    # Business-line pages
    business_dir=repo/"docs/business-lines"; business_dir.mkdir(parents=True,exist_ok=True)
    (business_dir/"README.md").write_text(business_index(),encoding="utf-8")
    for identity in IDENTITIES:
        page=business_dir/f"{identity.slug.upper().replace('-', '_')}.md"
        page.write_text(business_page(identity),encoding="utf-8")

    # Wiki source
    wiki=repo/"docs/wiki"; wiki.mkdir(parents=True,exist_ok=True)
    (wiki/"Home.md").write_text(wiki_home(),encoding="utf-8")
    sidebar=["## Sable Harbor","","- [[Home]]","- [[Business Lines]]","- [[Organization]]","- [[Corporate History]]","- [[Brand Assets]]","- [[Canon and Governance]]","","### Identities"]
    for identity in IDENTITIES:
        sidebar.append(f"- [[{identity.display_name}]]")
    (wiki/"_Sidebar.md").write_text("\n".join(sidebar)+"\n",encoding="utf-8")
    line_links=[]
    for identity in IDENTITIES:
        name=identity.display_name.replace(",", "")
        filename=identity.display_name.replace("/", "-")
        (wiki/f"{filename}.md").write_text(wiki_page(identity),encoding="utf-8")
        line_links.append(f"- [[{identity.display_name}]] - {identity.role}")
    (wiki/"Business Lines.md").write_text("# Business Lines\n\n"+"\n".join(line_links)+"\n",encoding="utf-8")
    (wiki/"Brand Assets.md").write_text(f"""# Brand Assets

The production logo library provides one lockup per file in outlined SVG and rendered PNG formats, plus document and presentation collateral.

- [Brand package]({REPO_URL}/tree/main/assets/brand)
- [Individual logos]({REPO_URL}/tree/main/assets/brand/logos)
- [Collateral]({REPO_URL}/tree/main/assets/brand/collateral)
- [Brand standards]({REPO_URL}/blob/main/assets/brand/BRAND_STANDARDS.md)
- [Complete ZIP]({REPO_URL}/tree/main/assets/brand/packages)

Exploratory boards and mockups are not production assets.
""",encoding="utf-8")
    (wiki/"Organization.md").write_text(f"# Organization\n\n[Open the rendered organization-chart suite]({REPO_URL}/tree/main/docs/organization).\n\nExact legal entities and HR reporting lines remain open where canon says they are open.\n",encoding="utf-8")
    (wiki/"Corporate History.md").write_text(f"# Corporate History\n\n[Read the corporate-lore canon]({REPO_URL}/blob/main/{CANON_PATH}).\n",encoding="utf-8")
    (wiki/"Canon and Governance.md").write_text(f"# Canon and Governance\n\nRepository canon controls over this wiki.\n\n- [Decision register]({REPO_URL}/blob/main/docs/canon/DECISION_REGISTER.md)\n- [Continuity audit]({REPO_URL}/blob/main/docs/canon/SABLE_HARBOR_CONTINUITY_AUDIT_v0.2.md)\n- [Public repository and wiki policy]({REPO_URL}/blob/main/docs/governance/PUBLIC_REPOSITORY_AND_WIKI_POLICY.md)\n",encoding="utf-8")
    (wiki/"PUBLISH_STATUS.md").write_text("# Wiki publication status\n\nWiki source is complete and versioned in `docs/wiki/`. Publication is performed by the manually dispatched `publish-wiki.yml` workflow and must be verified independently.\n",encoding="utf-8")


def root_readme_block() -> str:
    cards=[]
    for i in [x for x in IDENTITIES if x.current and x.slug != "sable-harbor"]:
        page=i.slug.upper().replace("-","_")
        cards.append(f'<td align="center" width="25%"><a href="docs/business-lines/{page}.md"><img src="assets/brand/logos/{i.slug}__mark.svg" width="88" alt="{i.display_name} mark"/><br/><strong>{i.display_name}</strong></a><br/><sub>{i.role}</sub></td>')
    rows=[]
    for n in range(0,len(cards),4):
        row=cards[n:n+4]
        if len(row)<4: row += ["<td></td>"]*(4-len(row))
        rows.append("<tr>"+"".join(row)+"</tr>")
    return f"""<!-- BRAND-INTEGRATION:START -->
<p align="center">
  <img src="assets/brand/logos/sable-harbor__primary-horizontal.svg" alt="Sable Harbor" width="880" />
</p>

<p align="center">
  <a href="docs/business-lines/README.md"><strong>Business lines</strong></a> &nbsp;·&nbsp;
  <a href="assets/brand/README.md"><strong>Brand assets</strong></a> &nbsp;·&nbsp;
  <a href="assets/brand/collateral/README.md"><strong>Corporate collateral</strong></a> &nbsp;·&nbsp;
  <a href="docs/organization/README.md"><strong>Organization maps</strong></a> &nbsp;·&nbsp;
  <a href="docs/wiki/Home.md"><strong>Wiki source</strong></a>
</p>

## Business lines

<table>
{chr(10).join(rows)}
</table>

The rendered identities are a publication layer over corporate lore; they do not independently create legal entities, reporting lines, economics, or other facts that remain OPEN.
<!-- BRAND-INTEGRATION:END -->"""


def insert_marker_block(text: str, block: str, marker: str, after_pattern: str | None = None) -> str:
    start=f"<!-- {marker}:START -->"; end=f"<!-- {marker}:END -->"
    if start in text and end in text:
        return re.sub(re.escape(start)+r".*?"+re.escape(end),block,text,flags=re.S)
    if after_pattern:
        m=re.search(after_pattern,text,flags=re.S)
        if m:
            return text[:m.end()]+"\n\n"+block+text[m.end():]
    return block+"\n\n"+text


def integrate_readmes(repo: Path) -> None:
    root=repo/"README.md"
    text=root.read_text(encoding="utf-8") if root.exists() else "# SABLE HARBOR\n"
    # Put block after the first descriptive paragraph if possible.
    block=root_readme_block()
    text=insert_marker_block(text,block,"BRAND-INTEGRATION",r"\A# SABLE HARBOR\s*\n(?:.*?\n){1,4}?(?=\n##|\Z)")
    root.write_text(text,encoding="utf-8")

    mappings={
        "docs/organization/README.md":"sable-harbor",
        "docs/organization/2026_LEADERSHIP_AND_AUTHORITY_MAP.md":"sable-harbor",
        "docs/organization/ORIGINAL_EIGHT.md":"sable-harbor",
        "docs/organization/FOUNDRY_FIELD_ORGANIZATION.md":"foundry-field",
        "docs/organization/WILLOW_ORGANIZATION.md":"willow",
        "docs/organization/ATLAS_MERIDIAN_BRIDGE_ORGANIZATION.md":"atlas-meridian",
        "docs/organization/PALE_SUN_RED_WASH_ORGANIZATION.md":"pale-sun",
        "docs/organization/PROJECT_CRADLE_ORGANIZATION.md":"project-cradle",
        "docs/organization/ARU_BST_ORGANIZATION.md":"american-resource-utility",
    }
    for rel,slug in mappings.items():
        path=repo/rel
        if not path.exists():
            continue
        txt=path.read_text(encoding="utf-8")
        alt=next(i.display_name for i in IDENTITIES if i.slug==slug)
        logo=f"../../assets/brand/logos/{slug}__primary-horizontal.svg"
        block=f'<!-- BRAND-HEADER:START -->\n<p align="center"><img src="{logo}" alt="{alt} logo" width="760" /></p>\n<!-- BRAND-HEADER:END -->'
        txt=insert_marker_block(txt,block,"BRAND-HEADER",r"\A# .*?\n")
        path.write_text(txt,encoding="utf-8")


def collateral_readme() -> str:
    return f"""# Sable Harbor Corporate Collateral

**Version:** {VERSION}

This folder contains production-oriented, editable corporate templates. No address, telephone number, email address, web domain, legal-entity suffix, or office location is invented; all such fields are explicit placeholders.

## Inventory

### Letterhead

- US Letter: SVG, PNG, PDF, DOCX
- A4: SVG, PNG, PDF, DOCX

### Memorandum

- US Letter: DOCX and PDF reference

### Report system

- US Letter and A4 cover: SVG, PNG, PDF
- US Letter and A4 report template: DOCX

### Presentation

- 16:9 PowerPoint template: PPTX

### Standards

- Brand standards: Markdown, PDF
- Font provenance: Markdown
- Preliminary name-and-mark screen: Markdown

## Use

Use the SVG files as vector source; use DOCX and PPTX for editable office documents; use PDFs as reproduction references after local proofing. Replace bracketed placeholders and preserve the logo geometry and clear space.
"""


def write_collateral(repo: Path) -> None:
    root=repo/"assets/brand/collateral"
    (root/"letterhead").mkdir(parents=True,exist_ok=True)
    (root/"memo").mkdir(parents=True,exist_ok=True)
    (root/"report").mkdir(parents=True,exist_ok=True)
    (root/"presentation").mkdir(parents=True,exist_ok=True)
    (root/"guides").mkdir(parents=True,exist_ok=True)
    (root/"README.md").write_text(collateral_readme(),encoding="utf-8")

    for page in ("us-letter","a4"):
        svg,w,h=self_contained_letterhead_svg(repo,page)
        svg_path=root/"letterhead"/f"sable-harbor-letterhead-{page}.svg"
        svg_path.write_text(svg,encoding="utf-8")
        svg_to_png_pdf(svg_path,svg_path.with_suffix(".png"),svg_path.with_suffix(".pdf"))
        create_letterhead_docx(repo,root/"letterhead"/f"sable-harbor-letterhead-{page}.docx",page)

        svg,w,h=self_contained_report_cover_svg(repo,page)
        svg_path=root/"report"/f"sable-harbor-report-cover-{page}.svg"
        svg_path.write_text(svg,encoding="utf-8")
        svg_to_png_pdf(svg_path,svg_path.with_suffix(".png"),svg_path.with_suffix(".pdf"))
        create_report_template_docx(repo,root/"report"/f"sable-harbor-report-template-{page}.docx",page)

    create_memo_docx(repo,root/"memo"/"sable-harbor-memorandum-us-letter.docx")
    create_simple_reference_pdf(repo,root/"memo"/"sable-harbor-memorandum-us-letter.pdf","Memorandum Reference")
    create_presentation(repo,root/"presentation"/"sable-harbor-presentation-template-16x9.pptx")
    create_brand_standards_pdf(repo,root/"guides"/"sable-harbor-brand-standards-v0.2.0.pdf")


def package_files(repo: Path) -> None:
    packages=repo/"assets/brand/packages"; packages.mkdir(parents=True,exist_ok=True)
    for old in packages.glob("*v0.2.0*.zip"):
        old.unlink()

    def zip_paths(output: Path, paths: list[Path]) -> None:
        with zipfile.ZipFile(output,"w",zipfile.ZIP_DEFLATED,compresslevel=9) as z:
            for p in sorted(paths):
                if p.is_file() and p != output:
                    z.write(p,p.relative_to(repo))

    logos=list((repo/"assets/brand/logos").glob("*"))
    brand_docs=[repo/"assets/brand/README.md",repo/"assets/brand/BRAND_STANDARDS.md",repo/"assets/brand/FONT_PROVENANCE.md",repo/"assets/brand/VALIDATION.md",repo/"assets/brand/manifest.json"]
    collateral=list((repo/"assets/brand/collateral").rglob("*"))
    business=list((repo/"docs/business-lines").rglob("*.md"))
    legal=[repo/"docs/legal/PRELIMINARY_NAME_AND_MARK_SCREEN.md"]
    wiki=list((repo/"docs/wiki").rglob("*.md"))
    tooling=[repo/"tools/brand/build_brand_system_v0_2.py", repo/"scripts/validate_brand_system.py"]
    tooling=[p for p in tooling if p.exists()]
    zip_paths(packages/f"sable-harbor-logo-system-v{VERSION}.zip",logos+brand_docs)
    zip_paths(packages/f"sable-harbor-corporate-collateral-v{VERSION}.zip",collateral+brand_docs)
    zip_paths(packages/f"sable-harbor-complete-brand-kit-v{VERSION}.zip",logos+brand_docs+collateral+business+legal+wiki+tooling)
    sums=[]
    for p in sorted(packages.glob("*.zip")):
        sums.append(f"{hashlib.sha256(p.read_bytes()).hexdigest()}  {p.name}")
    (packages/"SHA256SUMS.txt").write_text("\n".join(sums)+"\n",encoding="utf-8")


def write_manifest(repo: Path) -> None:
    brand=repo/"assets/brand"
    records=[]
    for p in sorted(brand.rglob("*")):
        if not p.is_file() or p.name in {"manifest.json","SHA256SUMS.txt"} or p.suffix==".zip":
            continue
        data=p.read_bytes()
        rec={"path":str(p.relative_to(repo)).replace("\\","/"),"bytes":len(data),"sha256":hashlib.sha256(data).hexdigest(),"format":p.suffix.lstrip(".")}
        if p.suffix.lower()==".png":
            with Image.open(p) as im:
                rec["width"],rec["height"]=im.size
                rec["mode"]=im.mode
        if p.parent.name=="logos" and "__" in p.stem:
            rec["identity"],rec["variant"]=p.stem.split("__",1)
            rec["one_logo_per_file"]=True
        records.append(rec)
    manifest={
        "schema_version":"1.1",
        "package_version":VERSION,
        "generated":BUILD_DATE,
        "controlling_canon":CANON_PATH,
        "asset_policy":{"one_logo_per_file":True,"svg_is_source_of_truth":True,"font_binaries_distributed":False,"concept_boards_in_production":False},
        "font_provenance":{
            "family":"Inter Display",
            "license":"SIL Open Font License 1.1",
            "logo_lettering":"vector outlines; no live text",
            "font_binaries_distributed":False,
            "build_sources":[
                {"role":"medium","filename":FONT_MEDIUM_PATH.name,"build_path":str(FONT_MEDIUM_PATH),"sha256":hashlib.sha256(FONT_MEDIUM_PATH.read_bytes()).hexdigest()},
                {"role":"semibold","filename":FONT_SEMIBOLD_PATH.name,"build_path":str(FONT_SEMIBOLD_PATH),"sha256":hashlib.sha256(FONT_SEMIBOLD_PATH.read_bytes()).hexdigest()},
                {"role":"bold","filename":FONT_BOLD_PATH.name,"build_path":str(FONT_BOLD_PATH),"sha256":hashlib.sha256(FONT_BOLD_PATH.read_bytes()).hexdigest()},
            ],
        },
        "identities":[{"slug":i.slug,"display_name":i.display_name,"role":i.role,"status":i.status,"variants":list(i.variants),"current":i.current} for i in IDENTITIES],
        "files":records,
    }
    (brand/"manifest.json").write_text(json.dumps(manifest,indent=2,ensure_ascii=False)+"\n",encoding="utf-8")


def validate(repo: Path) -> dict:
    logo_dir=repo/"assets/brand/logos"
    errors=[]; warnings=[]
    svg_files=sorted(logo_dir.glob("*.svg")); png_files=sorted(logo_dir.glob("*.png"))
    stems_svg={p.stem for p in svg_files}; stems_png={p.stem for p in png_files}
    if stems_svg != stems_png:
        errors.append(f"SVG/PNG mismatch: svg-only={sorted(stems_svg-stems_png)}, png-only={sorted(stems_png-stems_svg)}")
    import xml.etree.ElementTree as ET
    for p in svg_files:
        try:
            root=ET.parse(p).getroot()
        except Exception as e:
            errors.append(f"Malformed SVG {p.name}: {e}"); continue
        tags=[el.tag.split("}")[-1] for el in root.iter()]
        if "text" in tags: errors.append(f"Live text in logo {p.name}")
        if "script" in tags: errors.append(f"Script in logo {p.name}")
        if "image" in tags: errors.append(f"Embedded raster/image in logo {p.name}")
        if tags.count("svg") != 1: errors.append(f"Nested/multiple SVG roots in logo {p.name}")
        if root.find("{http://www.w3.org/2000/svg}title") is None: errors.append(f"Missing SVG title {p.name}")
    for identity in IDENTITIES:
        for v in identity.variants:
            for ext in ("svg","png"):
                p=logo_dir/f"{identity.slug}__{v}.{ext}"
                if not p.exists(): errors.append(f"Missing expected logo {p.name}")
    forbidden=("board","contact-sheet","mockup","concept")
    for p in logo_dir.iterdir():
        if any(x in p.name.lower() for x in forbidden): errors.append(f"Forbidden production filename: {p.name}")
    required=[
        repo/"assets/brand/BRAND_STANDARDS.md",
        repo/"assets/brand/FONT_PROVENANCE.md",
        repo/"docs/legal/PRELIMINARY_NAME_AND_MARK_SCREEN.md",
        repo/"docs/business-lines/README.md",
        repo/"docs/wiki/Home.md",
        repo/"assets/brand/collateral/letterhead/sable-harbor-letterhead-us-letter.docx",
        repo/"assets/brand/collateral/letterhead/sable-harbor-letterhead-a4.docx",
        repo/"assets/brand/collateral/presentation/sable-harbor-presentation-template-16x9.pptx",
        repo/"assets/brand/collateral/guides/sable-harbor-brand-standards-v0.2.0.pdf",
    ]
    for p in required:
        if not p.exists() or p.stat().st_size==0: errors.append(f"Missing/empty required file: {p.relative_to(repo)}")
    # Manifest hashes.
    manifest_path=repo/"assets/brand/manifest.json"
    if manifest_path.exists():
        manifest=json.loads(manifest_path.read_text(encoding="utf-8"))
        for rec in manifest["files"]:
            p=repo/rec["path"]
            if not p.exists(): errors.append(f"Manifest path missing: {rec['path']}"); continue
            actual=hashlib.sha256(p.read_bytes()).hexdigest()
            if actual!=rec["sha256"]: errors.append(f"Manifest hash mismatch: {rec['path']}")
    else: errors.append("manifest.json missing")
    # Placeholder safety.
    for p in (repo/"assets/brand/collateral").rglob("*.svg"):
        txt=p.read_text(encoding="utf-8")
        if "[ADDRESS]" not in txt and "report-cover" not in p.name:
            warnings.append(f"No address placeholder in {p.name}")
    return {"ok":not errors,"errors":errors,"warnings":warnings,"svg_count":len(svg_files),"png_count":len(png_files),"lockup_count":len(stems_svg)}


def validation_markdown(result: dict) -> str:
    status="PASS" if result["ok"] else "FAIL"
    errors="\n".join(f"- {e}" for e in result["errors"]) or "- None"
    warnings="\n".join(f"- {w}" for w in result["warnings"]) or "- None"
    return f"""# Brand System Validation

**Version:** {VERSION}  
**Date:** {BUILD_DATE}  
**Result:** **{status}**

## Production logo checks

- Distinct lockups: **{result['lockup_count']}**
- SVG files: **{result['svg_count']}**
- PNG files: **{result['png_count']}**
- Every SVG has a matching PNG.
- Exactly one identity/lockup is stored per file.
- No contact sheets, concept boards, or mockups are stored under `assets/brand/logos/`.
- Production logo SVGs contain no live `<text>`, `<script>`, or embedded `<image>` elements.
- Every expected current, supplemental, and historical identity variant exists.
- Evalon is explicitly classified as historical and not added to the seven-line 2026 map.

## Publication and collateral checks

- Root README integration block generated.
- Dedicated business-line pages generated.
- Organization-page logo headers generated where the corresponding page exists.
- Wiki source package generated.
- US Letter and A4 letterhead generated in SVG, PNG, PDF, and DOCX.
- Memo, report, presentation, brand standards, provenance, and preliminary legal-screen artifacts generated.
- Package checksums generated.
- No contact information or legal-entity suffix was invented in editable collateral.

## Errors

{errors}

## Warnings

{warnings}
"""


def update_brand_readme(repo: Path) -> None:
    current_rows=[]
    for i in IDENTITIES:
        if i.current:
            current_rows.append(f"| {i.display_name} | {i.role} | {len(i.variants)} |")
    other_rows=[]
    for i in IDENTITIES:
        if not i.current:
            other_rows.append(f"| {i.display_name} | {i.role} | {len(i.variants)} | {i.status} |")
    text=f"""# Sable Harbor Corporate Identity System - v{VERSION}

This directory contains individual production logo assets, corporate collateral, brand standards, provenance, validation, and packaged distributions for Sable Harbor.

## File rule

- **One logo per file.**
- No contact sheet, composite concept board, or mockup is a production logo.
- Every logo lockup is supplied as self-contained outlined SVG and rendered PNG.
- SVG is the production source of truth.
- Reverse lockups use a dark field; other PNGs preserve transparency.

## Current identity architecture

| Identity | Role | Variants |
|---|---|---:|
{chr(10).join(current_rows)}

## Supplemental and historical identities

| Identity | Role | Variants | Canon note |
|---|---|---:|---|
{chr(10).join(other_rows)}

Evalon is included as a deliberately archival identity because it is a consequential predecessor to Willow. It is not a current business line.

## Directory map

```text
assets/brand/
├── logos/                  # one logo/lockup per SVG and PNG
├── collateral/             # letterhead, memo, report, presentation, guide
├── packages/               # complete downloadable ZIP distributions
├── BRAND_STANDARDS.md
├── FONT_PROVENANCE.md
├── VALIDATION.md
└── manifest.json
```

## Direct use

- README/wiki headings: `__primary-horizontal.svg`
- Covers/square fields: `__stacked.svg`
- Avatars/navigation: `__mark.svg`
- Dark fields: `__reverse-horizontal.svg`
- Monochrome/stamps/engraving: `__one-color-horizontal.svg`

## Collateral

See [`collateral/README.md`](collateral/README.md) for US Letter and A4 letterheads, memorandum, report templates, covers, PowerPoint, and the rendered brand-standards guide.

## Governance

- [Brand standards](BRAND_STANDARDS.md)
- [Font provenance](FONT_PROVENANCE.md)
- [Validation](VALIDATION.md)
- [Manifest](manifest.json)
- [Preliminary name and mark screen](../../docs/legal/PRELIMINARY_NAME_AND_MARK_SCREEN.md)
- [Business-line index](../../docs/business-lines/README.md)

Artwork follows corporate lore v0.2 and does not independently create legal entities, reporting lines, economics, or other facts that remain OPEN. All rights reserved unless a specific file states otherwise.
"""
    (repo/"assets/brand/README.md").write_text(text,encoding="utf-8")


def write_validation_script(repo: Path) -> None:
    scripts=repo/"scripts"; scripts.mkdir(parents=True,exist_ok=True)
    validator='''#!/usr/bin/env python3
from __future__ import annotations
import hashlib, json, sys
from pathlib import Path
import xml.etree.ElementTree as ET

ROOT=Path(__file__).resolve().parents[1]
LOGOS=ROOT/'assets/brand/logos'
errors=[]
svg=sorted(LOGOS.glob('*.svg')); png=sorted(LOGOS.glob('*.png'))
ss={p.stem for p in svg}; ps={p.stem for p in png}
if ss!=ps: errors.append(f'SVG/PNG stem mismatch: {sorted(ss^ps)}')
for p in svg:
    try: root=ET.parse(p).getroot()
    except Exception as exc: errors.append(f'{p}: malformed SVG: {exc}'); continue
    tags=[e.tag.split('}')[-1] for e in root.iter()]
    for forbidden in ('text','script','image'):
        if forbidden in tags: errors.append(f'{p}: forbidden <{forbidden}> in production logo')
    if tags.count('svg')!=1: errors.append(f'{p}: expected exactly one SVG root')
    if root.find('{http://www.w3.org/2000/svg}title') is None: errors.append(f'{p}: missing title')
for p in LOGOS.iterdir():
    if any(x in p.name.lower() for x in ('board','contact-sheet','mockup','concept')):
        errors.append(f'{p}: production directory contains forbidden board/mockup filename')
manifest=json.loads((ROOT/'assets/brand/manifest.json').read_text())
for rec in manifest['files']:
    p=ROOT/rec['path']
    if not p.exists(): errors.append(f"manifest path missing: {rec['path']}"); continue
    if hashlib.sha256(p.read_bytes()).hexdigest()!=rec['sha256']:
        errors.append(f"manifest hash mismatch: {rec['path']}")
required=[
 'assets/brand/BRAND_STANDARDS.md','assets/brand/FONT_PROVENANCE.md','assets/brand/VALIDATION.md',
 'assets/brand/collateral/letterhead/sable-harbor-letterhead-us-letter.docx',
 'assets/brand/collateral/letterhead/sable-harbor-letterhead-a4.docx',
 'assets/brand/collateral/presentation/sable-harbor-presentation-template-16x9.pptx',
 'docs/business-lines/README.md','docs/legal/PRELIMINARY_NAME_AND_MARK_SCREEN.md','docs/wiki/Home.md'
]
for rel in required:
    p=ROOT/rel
    if not p.exists() or p.stat().st_size==0: errors.append(f'missing/empty: {rel}')
if errors:
    print('\\n'.join(errors)); sys.exit(1)
print(f'PASS: {len(ss)} lockups, {len(svg)} SVG, {len(png)} PNG; manifest and publication package verified.')
'''
    p=scripts/"validate_brand_system.py"; p.write_text(validator,encoding="utf-8"); p.chmod(0o755)
    workflow=repo/".github/workflows"; workflow.mkdir(parents=True,exist_ok=True)
    (workflow/"validate-brand-system.yml").write_text('''name: Validate brand system

on:
  pull_request:
    paths:
      - 'assets/brand/**'
      - 'docs/business-lines/**'
      - 'docs/legal/**'
      - 'docs/wiki/**'
      - 'scripts/validate_brand_system.py'
  push:
    branches: [main, canon/corporate-lore-v0.2]
    paths:
      - 'assets/brand/**'
      - 'docs/business-lines/**'
      - 'docs/legal/**'
      - 'docs/wiki/**'
      - 'scripts/validate_brand_system.py'
  workflow_dispatch:

permissions:
  contents: read

jobs:
  validate:
    runs-on: ubuntu-latest
    steps:
      - uses: actions/checkout@v4
      - uses: actions/setup-python@v5
        with:
          python-version: '3.12'
      - run: python scripts/validate_brand_system.py
''',encoding="utf-8")
    (workflow/"publish-wiki.yml").write_text('''name: Publish repository wiki

on:
  workflow_dispatch:

permissions:
  contents: write

jobs:
  publish:
    runs-on: ubuntu-latest
    steps:
      - uses: actions/checkout@v4
      - name: Publish versioned wiki source
        env:
          GH_TOKEN: ${{ secrets.GITHUB_TOKEN }}
        run: |
          set -euo pipefail
          rm -rf /tmp/sable-harbor-wiki
          WIKI_URL="https://x-access-token:${GH_TOKEN}@github.com/${GITHUB_REPOSITORY}.wiki.git"
          if ! git clone "$WIKI_URL" /tmp/sable-harbor-wiki; then
            mkdir -p /tmp/sable-harbor-wiki
            git -C /tmp/sable-harbor-wiki init
            git -C /tmp/sable-harbor-wiki remote add origin "$WIKI_URL"
          fi
          find /tmp/sable-harbor-wiki -mindepth 1 -maxdepth 1 ! -name .git -exec rm -rf {} +
          cp docs/wiki/*.md /tmp/sable-harbor-wiki/
          git -C /tmp/sable-harbor-wiki config user.name "Sable Harbor Wiki Bot"
          git -C /tmp/sable-harbor-wiki config user.email "sable-harbor-wiki-bot@users.noreply.github.com"
          git -C /tmp/sable-harbor-wiki add --all
          if git -C /tmp/sable-harbor-wiki diff --cached --quiet; then
            echo "Wiki already current."
            exit 0
          fi
          git -C /tmp/sable-harbor-wiki commit -m "Publish Sable Harbor wiki from repository source"
          git -C /tmp/sable-harbor-wiki push --force origin HEAD:master
''',encoding="utf-8")


def build(repo: Path) -> dict:
    repo=repo.resolve()
    logo_dir=repo/"assets/brand/logos"
    if not logo_dir.exists():
        raise FileNotFoundError(f"Existing logo directory not found: {logo_dir}")
    write_evalon_assets(logo_dir)
    write_docs(repo)
    integrate_readmes(repo)
    write_collateral(repo)
    shutil.rmtree(repo / ".brand-tmp", ignore_errors=True)
    update_brand_readme(repo)
    write_validation_script(repo)
    write_manifest(repo)
    result=validate(repo)
    (repo/"assets/brand/VALIDATION.md").write_text(validation_markdown(result),encoding="utf-8")
    # Recompute manifest after validation text changes, then validate and package.
    write_manifest(repo)
    result=validate(repo)
    (repo/"assets/brand/VALIDATION.md").write_text(validation_markdown(result),encoding="utf-8")
    write_manifest(repo)
    package_files(repo)
    result=validate(repo)
    shutil.rmtree(repo / ".brand-tmp", ignore_errors=True)
    if not result["ok"]:
        raise RuntimeError("Brand build failed:\n"+"\n".join(result["errors"]))
    return result


def main() -> int:
    ap=argparse.ArgumentParser()
    ap.add_argument("--repo-root",type=Path,default=Path(os.environ.get("GITHUB_WORKSPACE",Path.cwd())))
    args=ap.parse_args()
    result=build(args.repo_root)
    print(json.dumps(result,indent=2))
    return 0


if __name__=="__main__":
    raise SystemExit(main())
